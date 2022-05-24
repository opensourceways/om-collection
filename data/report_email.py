#!/usr/bin/env python3
# -*- coding: utf-8 -*-
#
# Copyright 2020 The community Authors.
# A-Tune is licensed under the Mulan PSL v2.
# You can use this software according to the terms and conditions of the Mulan PSL v2.
# You may obtain a copy of Mulan PSL v2 at:
#     http://license.coscl.org.cn/MulanPSL2
# THIS SOFTWARE IS PROVIDED ON AN "AS IS" BASIS, WITHOUT WARRANTIES OF ANY KIND, EITHER EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO NON-INFRINGEMENT, MERCHANTABILITY OR FIT FOR A PARTICULAR
# PURPOSE.
# See the Mulan PSL v2 for more details.
# Create: 2020-05
#
import bs4
import os
import time
import datetime
from dateutil import tz
import json
import requests
import configparser
import glob
import threading
import requests
import pypistats
import re
from pprint import pprint
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from apscheduler.schedulers.blocking import BlockingScheduler
from data.sendEmail import send_email
from pptx.dml.color import RGBColor
from data import common
from data.common import ESClient


class ReportEmail(object):

    def __init__(self, config=None):
        self.config = config
        self.index_name = config.get('index_name')
        self.org = config.get('org')
        self.esClient = ESClient(config)
        self.github_authorization = config.get('github_authorization')
        self.url = config.get('es_url')
        self.from_data = config.get("from_data")
        self.headers = {'Content-Type': 'application/json', 'Authorization': config.get('authorization')}
        self.email_list = config.get('email_list').split(',')
        self.send = config.get('send')
        self.smtp = config.get('smtp')
        self.repo = config.get('repo')
        self.github_star = config.get('github_star')
        self.gitee_issue = config.get('gitee_issue')
        self.committersNotHuawei = config.get('committersnothuawei')
        self.github_watch = config.get('github_watch')
        self.gitee_fork_user = config.get('gitee_fork_user')
        self.contributorNotHuawei = config.get('contributornothuawei')
        self.downloadall = config.get('downloadall')
        self.contributor = config.get('contributor')
        self.downlaod_ip = config.get('downlaod_ip')
        self.contributorHuawei = config.get('contributorhuawei')
        self.maillist = config.get('maillist')
        self.gitee_star_user = config.get('gitee_star_user')
        self.D0_user = config.get('d0_user')
        self.downlaod = config.get('downlaod')
        self.deep_contributor = config.get('deep_contributor')
        self.webvisitor = config.get('webvisitor')
        self.users = config.get('users')
        self.committersHuawei = config.get('committershuawei')
        self.video = config.get('video')
        self.github_fork = config.get('github_fork')
        self.maillistnum = config.get('maillistnum')
        self.gitee_issue_user = config.get('gitee_issue_user')
        self.gitee_pr_user = config.get('gitee_pr_user')
        self.gitee_watch_user = config.get('gitee_watch_user')
        self.paly_video = config.get('paly_video')
        self.website = config.get('website')
        self.committers = config.get('committers')
        self.run_hour = config.get('run_hour')

    def run(self, from_date):
        
        scheduler = BlockingScheduler()
        scheduler.add_job(self.do, 'cron', hour=self.run_hour, kwargs={"org": self.org})
        scheduler.start()
        print('end...')

    def getPng(self, url, filename, api, start='2020-01-01 00:00:00', end=None):
        starttime = int(round(time.mktime(time.strptime(start, '%Y-%m-%d %H:%M:%S')) * 1000))
        if end:
            endtime = int(round(time.mktime(time.strptime(end, '%Y-%m-%d %H:%M:%S')) * 1000))
        else:
            endtime = int(round(time.time() * 1000))
        time.sleep(5)
        url = re.sub(r'from=\d+&to=\d+&panelId', 'from=%s&to=%s&panelId' % (starttime, endtime), url)
        print(url)
        cmd = 'curl -H "Accept: application/json" -H "Authorization: Bearer %s" "%s" > %s' % (api, url, filename)
        os.system(cmd)

    def getData(self, headers, start='2020-01-01 00:00:00', end=None, org=None, header01=''):
        result = {"openeuler":{}, 'mindspore':{}, 'opengauss':{}, 'openlookeng':{}}
        if end:
            endtime = end.replace(' ', 'T')
            LW_endtime = datetime.datetime.strftime(
                datetime.datetime.strptime(endtime, "%Y-%m-%dT%H:%M:%S") - datetime.timedelta(days=7), "%Y-%m-%dT%H:%M:%S")

        else:
            endtime = time.strftime("%Y-%m-%dT%H:%M:%S", time.localtime(time.time()))
            LW_endtime = datetime.datetime.strftime(
                datetime.datetime.strptime(endtime, "%Y-%m-%dT%H:%M:%S") - datetime.timedelta(days=7), "%Y-%m-%dT%H:%M:%S")

        starttime = start.replace(' ', 'T')
        # openeuler
        if org == 'openeuler':
            data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              },
                // "must": [{ "match": { "ip":"127.0.0.1" }}],
                "should": [{
                        "match_phrase": {
                            "source_type_title": "搜索引擎"
                        }
                    },
              {
                        "match_phrase": {
                            "source_type_title": "外部链接"
                        }
                    },
                    {
                        "match_phrase": {
                            "source_type_title": "直接访问"
                        }
                    }],
                    "minimum_should_match": 1
            }
          },"aggs": {
            "ip_count": {
              "sum": {
                "field": "ip_count"}}}
        }''' % (starttime, endtime)
            url = self.webvisitor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["webvisitor"] = result["openeuler"]["D0_user"] = int(res["aggregations"]["ip_count"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },
                    // "must": [{ "match": { "ip":"127.0.0.1" }}],
                    "should": [{
                            "match_phrase": {
                                "source_type_title": "搜索引擎"
                            }
                        },
                  {
                            "match_phrase": {
                                "source_type_title": "外部链接"
                            }
                        },
                        {
                            "match_phrase": {
                                "source_type_title": "直接访问"
                            }
                        }],
                        "minimum_should_match": 1
                }
              },"aggs": {
                "pv_count": {
                  "sum": {
                    "field": "pv_count"}}}
            }''' % (starttime, endtime)
            url = self.website
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["website"] = int(res["aggregations"]["pv_count"]["value"])

            data = '''{"size":100,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },
              "aggs": {
                "sum_value_gb": {
                  "sum": {
                    "field": "sum_value_gb"}}}
            }''' % (starttime, endtime)
            url = self.downlaod
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["downlaod"] = int(int(res["aggregations"]["sum_value_gb"]["value"]) / 6)
            result["openeuler"]["downlaod_gb"] = int(res["aggregations"]["sum_value_gb"]["value"])
            result["openeuler"]["paly_video"] = 'N/A'

            data = '''{"size":100,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          }
                        }
                      },
                      "aggs": {
                        "ip": {
                          "cardinality": {
                            "field": "ip.keyword"}}}
                    }''' % (starttime, endtime)
            url = self.downlaod_ip
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["downlaod_ip"] = int(res["aggregations"]["ip"]["value"])

            data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              }
            }
          },"aggs": {
            "user_id": {
              "cardinality": {
                "field": "user_id.keyword"}}}
        }''' % (starttime, endtime)
            url = self.maillist
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["maillist"] = int(res["aggregations"]["user_id"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "metadata__updated_on":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }, "must": [{ "match": {"project":"openeuler_mail" }}]
                }
              },"aggs": {
                "user_id": {
                  "sum": {
                    "field": "is_hyperkitty_message"}}}
            }''' % (starttime, endtime)
            url = self.maillistnum
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["maillistnum"] = int(res["aggregations"]["user_id"]["value"])

            data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              }, "must": [{ "match": { "event":"DOWNLOAD ZIP" }}]
            }
          },
          "aggs": {
            "author_name": {
              "cardinality": {
                "field": "author_name.keyword"}}}
        }''' % (starttime, endtime)
            url = self.repo
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["repo"] = int(res["aggregations"]["author_name"]["value"])

            data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              }, "must_not": [{ "match": { "is_topic":1 }}]
            }
          },
          "aggs": {
            "new_followers": {
              "sum": {
                "field": "new_followers"}}}
        }''' % (starttime, endtime)
            url = self.video
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["video"] = int(res["aggregations"]["new_followers"]["value"])

            data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              }
            }
          },
          "aggs": {
            "sum_value_gb": {
              "sum": {
                "field": "sum_value_gb"}}}
        }''' % (starttime, endtime)
            url = self.users
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["users"] = int(res["aggregations"]["sum_value_gb"]["value"] / 6)

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },
              "aggs": {
                "sum_value_gb": {
                  "sum": {
                    "field": "sum_value_gb"}}}
            }''' % (starttime, endtime)
            url = self.downloadall
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["downloadall"] = int(res["aggregations"]["sum_value_gb"]["value"] / 6)

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      }
                    }
                  },"aggs": {
                    "user_login": {
                      "sum": {
                        "field": "is_gitee_issue"}}}
                }''' % (starttime, endtime)
            url = self.gitee_issue
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["gitee_issue"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },"aggs": {
                "user_login": {
                  "sum": {
                    "field": "is_gitee_star"}}}
            }''' % (starttime, endtime)
            url = self.gitee_star_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["gitee_star_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      }
                    }
                  },"aggs": {
                    "user_login": {
                      "sum": {
                        "field": "is_gitee_fork"}}}
                }''' % (starttime, endtime)
            url = self.gitee_fork_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["gitee_fork_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      }
                    }
                  },"aggs": {
                    "user_login": {
                      "sum": {
                        "field": "is_gitee_watch"}}}
                }''' % (starttime, endtime)
            url = self.gitee_watch_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["gitee_watch_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          }
                        }
                      },"aggs": {
                        "user_login": {
                          "sum": {
                            "field": "is_gitee_pull_request"}}}
                    }''' % (starttime, endtime)
            url = self.gitee_pr_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["gitee_pr_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                          "query": {
                            "bool": {
                              "filter":{
                                "range":{
                                  "created_at":{
                                    "gte":"%s.000+0800",
                                    "lt":"%s.000+0800"
                                  }
                                }
                              }
                            }
                          },"aggs": {
                            "user_login": {
                              "sum": {
                                "field": "is_gitee_issue"}}}
                        }''' % (starttime, endtime)
            url = self.gitee_issue_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["gitee_issue_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{
              "size": 0,
              "query": {
                "bool": {
                  "filter": [
                    {
                      "range": {
                        "created_at": {
                          "gte": "%s.000+0800",
                          "lte": "%s.000+0800"
                        }
                      }
                    },
                    {
                      "query_string": {
                        "analyze_wildcard": true,
                        "query": "*"
                      }
                    }
                  ]
                }
              },
              "aggs": {
                "2": {
                  "terms": {
                    "field": "html_url.keyword",
                    "size": 500,
                    "order": {
                      "_key": "desc"
                    },
                    "min_doc_count": 1
                  },
                  "aggs": {
                    "3": {
                      "max": {
                        "field": "forks_count"
                      }
                    }
                  }
                }
              }
            }''' % (starttime, endtime)
            url = self.github_fork
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["github_fork"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{
                  "size": 0,
                  "query": {
                    "bool": {
                      "filter": [
                        {
                          "range": {
                            "created_at": {
                              "gte": "%s.000+0800",
                              "lte": "%s.000+0800"
                            }
                          }
                        },
                        {
                          "query_string": {
                            "analyze_wildcard": true,
                            "query": "*"
                          }
                        }
                      ]
                    }
                  },
                  "aggs": {
                    "2": {
                      "terms": {
                        "field": "html_url.keyword",
                        "size": 500,
                        "order": {
                          "_key": "desc"
                        },
                        "min_doc_count": 1
                      },
                      "aggs": {
                        "3": {
                          "max": {
                            "field": "subscribers_count"
                          }
                        }
                      }
                    }
                  }
                }''' % (starttime, endtime)
            url = self.github_watch
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["github_watch"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{
                      "size": 0,
                      "query": {
                        "bool": {
                          "filter": [
                            {
                              "range": {
                                "created_at": {
                                  "gte": "%s.000+0800",
                                  "lte": "%s.000+0800"
                                }
                              }
                            },
                            {
                              "query_string": {
                                "analyze_wildcard": true,
                                "query": "*"
                              }
                            }
                          ]
                        }
                      },
                      "aggs": {
                        "2": {
                          "terms": {
                            "field": "html_url.keyword",
                            "size": 500,
                            "order": {
                              "_key": "desc"
                            },
                            "min_doc_count": 1
                          },
                          "aggs": {
                            "3": {
                              "max": {
                                "field": "stargazers_count"
                              }
                            }
                          }
                        }
                      }
                    }''' % (starttime, endtime)
            url = self.github_star
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["github_star"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              },
                // "must": [{ "match": { "ip":"127.0.0.1" }}],
                "should": [{
                        "match_phrase": {
                            "is_gitee_fork": 1
                        }
                    },
              {
                        "match_phrase": {
                            "is_gitee_issue": 1
                        }
                    },
                    {
                        "match_phrase": {
                            "is_gitee_pull_request": 1
                        }
                    },
                    {
                        "match_phrase": {
                            "is_gitee_review_comment": 1
                        }
                    },
                    {
                        "match_phrase": {
                            "is_gitee_comment": 1
                        }
                    },
                    {
                        "match_phrase": {
                            "is_gitee_issue_comment": 1
                        }
                    }],
                    "minimum_should_match": 1
            }
          },"aggs": {
            "user_login": {
              "cardinality": {
                "field": "user_login.keyword"}}}
        }''' % (starttime, endtime)
            url = self.contributor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["contributor"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },
                    // "must": [{ "match": { "ip":"127.0.0.1" }}],
                    "should": [
                  {
                            "match_phrase": {
                                "is_gitee_issue": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_pull_request": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_review_comment": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_comment": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_issue_comment": 1
                            }
                        }],
                        "minimum_should_match": 1
                }
              },"aggs": {
                "user_login": {
                  "cardinality": {
                    "field": "user_login.keyword"}}}
            }''' % (starttime, endtime)
            url = self.deep_contributor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["deep_contributor"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              },
                "must": [{ "match": { "is_project_internal_user":1 }}],
                "should": [{
                        "match_phrase": {
                            "is_gitee_fork": 1
                        }
                    },
              {
                        "match_phrase": {
                            "is_gitee_issue": 1
                        }
                    },
                    {
                        "match_phrase": {
                            "is_gitee_pull_request": 1
                        }
                    },
                    {
                        "match_phrase": {
                            "is_gitee_review_comment": 1
                        }
                    },
                    {
                        "match_phrase": {
                            "is_gitee_comment": 1
                        }
                    },
                    {
                        "match_phrase": {
                            "is_gitee_issue_comment": 1
                        }
                    }],
                    "minimum_should_match": 1
            }
          },"aggs": {
            "user_login": {
              "cardinality": {
                "field": "user_login.keyword"}}}
        }''' % (starttime, endtime)
            url = self.contributorHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["contributorHuawei"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },
                    "must_not": [{ "match": { "is_project_internal_user":1 }}],
                    "should": [{
                            "match_phrase": {
                                "is_gitee_fork": 1
                            }
                        },
                  {
                            "match_phrase": {
                                "is_gitee_issue": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_pull_request": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_review_comment": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_comment": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_issue_comment": 1
                            }
                        }],
                        "minimum_should_match": 1
                }
              },"aggs": {
                "user_login": {
                  "cardinality": {
                    "field": "user_login.keyword"}}}
            }''' % (starttime, endtime)
            url = self.contributorNotHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["contributorNotHuawei"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              }
            }
          },"aggs": {
            "user_gitee_name": {
              "cardinality": {
                "field": "committer.keyword"}}}
        }''' % (starttime, endtime)
            url = self.committers
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["committers"] = int(res["aggregations"]["user_gitee_name"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },"must": [{ "match": { "is_project_internal_user":1 }}]
                }
              },"aggs": {
                "user_gitee_name": {
                  "cardinality": {
                    "field": "committer.keyword"}}}
            }''' % (starttime, endtime)
            url = self.committersHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["committersHuawei"] = int(res["aggregations"]["user_gitee_name"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },"must_not": [{ "match": { "is_project_internal_user":1 }}]
                }
              },"aggs": {
                "user_gitee_name": {
                  "cardinality": {
                    "field": "committer.keyword"}}}
            }''' % (starttime, endtime)
            url = self.committersNotHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openeuler"]["committersNotHuawei"] = int(res["aggregations"]["user_gitee_name"]["value"])

        elif org == 'mindspore':
            # mindspore
            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },
                    // "must": [{ "match": { "ip":"127.0.0.1" }}],
                    "should": [{
                            "match_phrase": {
                                "source_type_title": "搜索引擎"
                            }
                        },
                  {
                            "match_phrase": {
                                "source_type_title": "外部链接"
                            }
                        },
                        {
                            "match_phrase": {
                                "source_type_title": "直接访问"
                            }
                        }],
                        "minimum_should_match": 1
                }
              },"aggs": {
                "ip_count": {
                  "sum": {
                    "field": "ip_count"}}}
            }''' % (starttime, endtime)
            url = self.webvisitor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["webvisitor"] = result["mindspore"]["D0_user"] = int(res["aggregations"]["ip_count"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      },
                        // "must": [{ "match": { "ip":"127.0.0.1" }}],
                        "should": [{
                                "match_phrase": {
                                    "source_type_title": "搜索引擎"
                                }
                            },
                      {
                                "match_phrase": {
                                    "source_type_title": "外部链接"
                                }
                            },
                            {
                                "match_phrase": {
                                    "source_type_title": "直接访问"
                                }
                            }],
                            "minimum_should_match": 1
                    }
                  },"aggs": {
                    "pv_count": {
                      "sum": {
                        "field": "pv_count"}}}
                }''' % (starttime, endtime)
            url = self.website
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["website"] = int(res["aggregations"]["pv_count"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },"aggs": {
                "user_id": {
                  "cardinality": {
                    "field": "user_id.keyword"}}}
            }''' % (starttime, endtime)
            url = self.maillist
            res = json.loads(requests.get(url=url, headers=header01, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["maillist"] = int(res["aggregations"]["user_id"]["value"])

            data = '''{"size":10000,
                          "query": {
                            "bool": {
                              "filter":{
                                "range":{
                                  "metadata__updated_on":{
                                    "gte":"%s.000+0800",
                                    "lt":"%s.000+0800"
                                  }
                                }
                              }, "must": [{ "match": {"project":"mindspore-mail" }}]
                            }
                          },"aggs": {
                            "user_id": {
                              "sum": {
                                "field": "is_hyperkitty_message"}}}
                        }''' % (starttime, endtime)
            url = self.maillistnum
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["maillistnum"] = int(res["aggregations"]["user_id"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }, "must": [{ "match": { "event":"DOWNLOAD ZIP" }},{ "match": { "is_forked_repo":0 }}]
                }
              },
              "aggs": {
                "author_name": {
                  "cardinality": {
                    "field": "author_name.keyword"}}}
            }''' % (starttime, endtime)
            url = self.repo
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["repo"] = int(res["aggregations"]["author_name"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }, "must_not": [{ "match": { "is_topic":1 }}]
                }
              },
              "aggs": {
                "new_followers": {
                  "sum": {
                    "field": "new_followers"}}}
            }''' % (starttime, endtime)
            url = self.video
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["video"] = int(res["aggregations"]["new_followers"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },
              "aggs": {
                "ip": {
                  "cardinality": {
                    "field": "ip.keyword"}}}
            }''' % (starttime, endtime)
            url = self.users
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["users"] = int(res["aggregations"]["ip"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      },
                        "must": [{ "match": { "is_gitee_star":1 }}]
                    }
                  },"aggs": {
                    "user_login": {
                      "cardinality": {
                        "field": "user_login.keyword"}}}
                }''' % (starttime, endtime)
            url = self.gitee_star_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["gitee_star_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          },
                            "must": [{ "match": { "is_gitee_fork":1 }}]
                        }
                      },"aggs": {
                        "user_login": {
                          "cardinality": {
                            "field": "user_login.keyword"}}}
                    }''' % (starttime, endtime)
            url = self.gitee_fork_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["gitee_fork_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          },
                            "must": [{ "match": { "is_gitee_watch":1 }}]
                        }
                      },"aggs": {
                        "user_login": {
                          "cardinality": {
                            "field": "user_login.keyword"}}}
                    }''' % (starttime, endtime)
            url = self.gitee_watch_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["gitee_watch_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                          "query": {
                            "bool": {
                              "filter":{
                                "range":{
                                  "created_at":{
                                    "gte":"%s.000+0800",
                                    "lt":"%s.000+0800"
                                  }
                                }
                              },
                                "must": [{ "match": { "is_gitee_pull_request":1 }}]
                            }
                          },"aggs": {
                            "user_login": {
                              "cardinality": {
                                "field": "user_login.keyword"}}}
                        }''' % (starttime, endtime)
            url = self.gitee_pr_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["gitee_pr_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                                      "query": {
                                        "bool": {
                                          "filter":{
                                            "range":{
                                              "created_at":{
                                                "gte":"%s.000+0800",
                                                "lt":"%s.000+0800"
                                              }
                                            }
                                          }
                                        }
                                      },"aggs": {
                                        "user_login": {
                                          "sum": {
                                            "field": "is_gitee_issue"}}}
                                    }''' % (starttime, endtime)
            url = self.gitee_issue
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["gitee_issue"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{
                          "size": 0,
                          "query": {
                            "bool": {
                              "filter": [
                                {
                                  "range": {
                                    "created_at": {
                                      "gte": "%s.000+0800",
                                      "lte": "%s.000+0800"
                                    }
                                  }
                                },
                                {
                                  "query_string": {
                                    "analyze_wildcard": true,
                                    "query": "*"
                                  }
                                }
                              ]
                            }
                          },
                          "aggs": {
                            "2": {
                              "terms": {
                                "field": "html_url.keyword",
                                "size": 500,
                                "order": {
                                  "_key": "desc"
                                },
                                "min_doc_count": 1
                              },
                              "aggs": {
                                "3": {
                                  "max": {
                                    "field": "forks_count"
                                  }
                                }
                              }
                            }
                          }
                        }''' % (starttime, endtime)
            url = self.github_fork
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["github_fork"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{
                              "size": 0,
                              "query": {
                                "bool": {
                                  "filter": [
                                    {
                                      "range": {
                                        "created_at": {
                                          "gte": "%s.000+0800",
                                          "lte": "%s.000+0800"
                                        }
                                      }
                                    },
                                    {
                                      "query_string": {
                                        "analyze_wildcard": true,
                                        "query": "*"
                                      }
                                    }
                                  ]
                                }
                              },
                              "aggs": {
                                "2": {
                                  "terms": {
                                    "field": "html_url.keyword",
                                    "size": 500,
                                    "order": {
                                      "_key": "desc"
                                    },
                                    "min_doc_count": 1
                                  },
                                  "aggs": {
                                    "3": {
                                      "max": {
                                        "field": "subscribers_count"
                                      }
                                    }
                                  }
                                }
                              }
                            }''' % (starttime, endtime)
            url = self.github_watch
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["github_watch"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{
                                  "size": 0,
                                  "query": {
                                    "bool": {
                                      "filter": [
                                        {
                                          "range": {
                                            "created_at": {
                                              "gte": "%s.000+0800",
                                              "lte": "%s.000+0800"
                                            }
                                          }
                                        },
                                        {
                                          "query_string": {
                                            "analyze_wildcard": true,
                                            "query": "*"
                                          }
                                        }
                                      ]
                                    }
                                  },
                                  "aggs": {
                                    "2": {
                                      "terms": {
                                        "field": "html_url.keyword",
                                        "size": 500,
                                        "order": {
                                          "_key": "desc"
                                        },
                                        "min_doc_count": 1
                                      },
                                      "aggs": {
                                        "3": {
                                          "max": {
                                            "field": "stargazers_count"
                                          }
                                        }
                                      }
                                    }
                                  }
                                }''' % (starttime, endtime)
            url = self.github_star
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["github_star"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at": {
                          "gte": "%s.000+0800",
                          "lte": "%s.000+0800"
                      }
                    }
                  }
                }
              }
            }''' % (starttime, endtime)
            url = self.paly_video
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["paly_video"] = int(res["count"])

            data = '''{
                        "query": {
                        "bool": {
                        "filter":{
                            "range":{
                              "created_at": {
                                  "gte": "%s.000+0800",
                                  "lte": "%s.000+0800"
                              }
                            }
                          },
                                "must": [{ "match": { "hostname.keyword":"download.mindspore.cn" }}]
                            }
                            }
                    }''' % (starttime, endtime)
            url = self.downlaod
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["downlaod"] = int(res["count"])
            result["mindspore"]["downlaod_gb"] = "N/A"

            data = '''{"size":10000,
                              "query": {
                                "bool": {
                                  "filter":{
                                    "range":{
                                      "created_at":{
                                        "gte":"%s.000+0800",
                                        "lt":"%s.000+0800"
                                      }
                                    }
                                  },"must": [{ "match": { "hostname":"download.mindspore.cn" }}]
                                }
                              },"aggs": {
                                "ip": {
                                  "cardinality": {
                                    "field": "ip.keyword"}}}
                            }''' % (starttime, endtime)
            url = self.downlaod_ip
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["downlaod_ip"] = int(res["aggregations"]["ip"]["value"])

            data = '''{"size":10000,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          },
                            "must_not": [{ "match": { "user_login":"mindspore_ci" }}],
                            "should": [{
                                    "match_phrase": {
                                        "is_gitee_fork": 1
                                    }
                                },
                          {
                                    "match_phrase": {
                                        "is_gitee_issue": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_pull_request": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_review_comment": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_comment": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_issue_comment": 1
                                    }
                                }],
                                "minimum_should_match": 1
                        }
                      },"aggs": {
                        "user_login": {
                          "cardinality": {
                            "field": "user_login.keyword"}}}
                    }''' % (starttime, endtime)
            url = self.contributor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["contributor"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                              "query": {
                                "bool": {
                                  "filter":{
                                    "range":{
                                      "created_at":{
                                        "gte":"%s.000+0800",
                                        "lt":"%s.000+0800"
                                      }
                                    }
                                  },
                                    "must_not": [{ "match": { "user_login":"mindspore_ci" }}],
                                    "should": [
                                  {
                                            "match_phrase": {
                                                "is_gitee_issue": 1
                                            }
                                        },
                                        {
                                            "match_phrase": {
                                                "is_gitee_pull_request": 1
                                            }
                                        },
                                        {
                                            "match_phrase": {
                                                "is_gitee_review_comment": 1
                                            }
                                        },
                                        {
                                            "match_phrase": {
                                                "is_gitee_comment": 1
                                            }
                                        },
                                        {
                                            "match_phrase": {
                                                "is_gitee_issue_comment": 1
                                            }
                                        }],
                                        "minimum_should_match": 1
                                }
                              },"aggs": {
                                "user_login": {
                                  "cardinality": {
                                    "field": "user_login.keyword"}}}
                            }''' % (starttime, endtime)
            url = self.deep_contributor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["deep_contributor"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },
                    "must": [{ "match": { "is_project_internal_user":1 }}],
                    "should": [{
                            "match_phrase": {
                                "is_gitee_fork": 1
                            }
                        },
                  {
                            "match_phrase": {
                                "is_gitee_issue": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_pull_request": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_review_comment": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_comment": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_issue_comment": 1
                            }
                        }],
                        "minimum_should_match": 1
                }
              },"aggs": {
                "user_login": {
                  "cardinality": {
                    "field": "user_login.keyword"}}}
            }''' % (starttime, endtime)
            url = self.contributorHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["contributorHuawei"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      },
                        "must": [{ "match": { "is_project_internal_user":0 }}],
                        "should": [{
                                "match_phrase": {
                                    "is_gitee_fork": 1
                                }
                            },
                      {
                                "match_phrase": {
                                    "is_gitee_issue": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_pull_request": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_review_comment": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_comment": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_issue_comment": 1
                                }
                            }],
                            "minimum_should_match": 1
                    }
                  },"aggs": {
                    "user_login": {
                      "cardinality": {
                        "field": "user_login.keyword"}}}
                }''' % (starttime, endtime)
            url = self.contributorNotHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["contributorNotHuawei"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },"must": [{ "match": { "is_committer":1 }}]
                }
              },"aggs": {
                "user_login": {
                  "cardinality": {
                    "field": "user_login.keyword"}}}
            }''' % (starttime, endtime)
            url = self.committers
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["mindspore"]["committers"] = int(res["aggregations"]["user_login"]["value"])

        elif org == 'opengauss':
            # opengauss
            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },
                    // "must": [{ "match": { "ip":"127.0.0.1" }}],
                    "should": [{
                            "match_phrase": {
                                "source_type_title": "搜索引擎"
                            }
                        },
                  {
                            "match_phrase": {
                                "source_type_title": "外部链接"
                            }
                        },
                        {
                            "match_phrase": {
                                "source_type_title": "直接访问"
                            }
                        }],
                        "minimum_should_match": 1
                }
              },"aggs": {
                "ip_count": {
                  "sum": {
                    "field": "ip_count"}}}
            }''' % (starttime, endtime)
            url = self.webvisitor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["webvisitor"] = int(res["aggregations"]["ip_count"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      },
                        // "must": [{ "match": { "ip":"127.0.0.1" }}],
                        "should": [{
                                "match_phrase": {
                                    "source_type_title": "搜索引擎"
                                }
                            },
                      {
                                "match_phrase": {
                                    "source_type_title": "外部链接"
                                }
                            },
                            {
                                "match_phrase": {
                                    "source_type_title": "直接访问"
                                }
                            }],
                            "minimum_should_match": 1
                    }
                  },"aggs": {
                    "pv_count": {
                      "sum": {
                        "field": "pv_count"}}}
                }''' % (starttime, endtime)
            url = self.website
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["website"] = int(res["aggregations"]["pv_count"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },"aggs": {
                "user_id": {
                  "cardinality": {
                    "field": "user_id.keyword"}}}
            }''' % (starttime, endtime)
            url = self.maillist
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["maillist"] = int(res["aggregations"]["user_id"]["value"])

            data = '''{"size":10000,
                          "query": {
                            "bool": {
                              "filter":{
                                "range":{
                                  "metadata__updated_on":{
                                    "gte":"%s.000+0800",
                                    "lt":"%s.000+0800"
                                  }
                                }
                              }, "must": [{ "match": {"project":"opengauss-mail" }}]
                            }
                          },"aggs": {
                            "user_id": {
                              "sum": {
                                "field": "is_hyperkitty_message"}}}
                        }''' % (starttime, endtime)
            url = self.maillistnum
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["maillistnum"] = int(res["aggregations"]["user_id"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }, "must": [{ "match": { "event":"DOWNLOAD ZIP" }},{ "match": { "is_forked_repo":1 }}]
                }
              },
              "aggs": {
                "author_name": {
                  "cardinality": {
                    "field": "author_name.keyword"}}}
            }''' % (starttime, endtime)
            url = self.repo
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["repo"] = int(res["aggregations"]["author_name"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }, "must_not": [{ "match": { "is_topic":1 }}]
                }
              },
              "aggs": {
                "new_followers": {
                  "sum": {
                    "field": "new_followers"}}}
            }''' % (starttime, endtime)
            url = self.video
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["video"] = int(res["aggregations"]["new_followers"]["value"])

            data = '''{
              "size": 0,
              "query": {
                "bool": {
                  "filter": [
                    {
                      "range": {
                        "created_at": {
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                        }
                      }
                    },
                    {
                      "query_string": {
                        "analyze_wildcard": true,
                        "query": "(path:\\"basicdb.*\\" OR path:\\"opengauss.*\\") AND !path:\\"odbc.tar.gz.sha256\\" AND !path:\\"64bit.tar.gz.sha256\\"  AND !path:\\"jdbc.tar.gz.sha256\\" AND !path:\\"symbol.tar.gz.sha256\\" AND !path:\\"libpq.tar.gz.sha256\\" AND !path.keyword:\\"1.0.0/openGauss-third_party_binarylibs.tar.gz/info/refs\\" AND !path.keyword:\\"1.0.0/openGauss-third\\""
                      }
                    }
                  ]
                }
              },
              "aggs": {
                "3": {
                  "terms": {
                    "field": "buket_name.keyword",
                    "size": 500,
                    "order": {
                      "_count": "desc"
                    },
                    "min_doc_count": 1
                  },
                  "aggs": {
                    "5": {
                      "cardinality": {
                        "field": "ip.keyword"
                      }
                    },
                    "6": {
                      "cardinality": {
                        "field": "city.keyword"
                      }
                    },
                    "8": {
                      "cardinality": {
                        "field": "country.keyword"
                      }
                    }
                  }
                }
              }
            }''' % (starttime, endtime)
            url = self.paly_video
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["paly_video"] = res["aggregations"]["3"]["buckets"][0]['doc_count']
            result["opengauss"]["downlaod"] = res["aggregations"]["3"]["buckets"][1]['doc_count']
            result["opengauss"]["downlaod_ip"] = res["aggregations"]["3"]["buckets"][1]['5']['value']
            result["opengauss"]["downlaod_gb"] = "N/A"

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },
              "aggs": {
                "ip": {
                  "cardinality": {
                    "field": "ip.keyword"}}}
            }''' % (starttime, endtime)
            url = self.users
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["users"] = int(res["aggregations"]["ip"]["value"])

            data = '''{
                  "size": 0,
                  "query": {
                    "bool": {
                      "filter": [
                        {
                          "range": {
                            "created_at": {
                              "gte": "%s.000+0800",
                              "lte": "%s.000+0800"
                            }
                          }
                        },
                        {
                          "query_string": {
                            "analyze_wildcard": true,
                            "query": "*"
                          }
                        }
                      ]
                    }
                  },
                  "aggs": {
                    "2": {
                      "terms": {
                        "field": "html_url.keyword",
                        "size": 500,
                        "order": {
                          "_key": "desc"
                        },
                        "min_doc_count": 1
                      },
                      "aggs": {
                        "3": {
                          "max": {
                            "field": "forks_count"
                          }
                        }
                      }
                    }
                  }
                }''' % (starttime, endtime)
            url = self.github_fork
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["github_fork"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{
                      "size": 0,
                      "query": {
                        "bool": {
                          "filter": [
                            {
                              "range": {
                                "created_at": {
                                  "gte": "%s.000+0800",
                                  "lte": "%s.000+0800"
                                }
                              }
                            },
                            {
                              "query_string": {
                                "analyze_wildcard": true,
                                "query": "*"
                              }
                            }
                          ]
                        }
                      },
                      "aggs": {
                        "2": {
                          "terms": {
                            "field": "html_url.keyword",
                            "size": 500,
                            "order": {
                              "_key": "desc"
                            },
                            "min_doc_count": 1
                          },
                          "aggs": {
                            "3": {
                              "max": {
                                "field": "subscribers_count"
                              }
                            }
                          }
                        }
                      }
                    }''' % (starttime, endtime)
            url = self.github_watch
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["github_watch"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{
                          "size": 0,
                          "query": {
                            "bool": {
                              "filter": [
                                {
                                  "range": {
                                    "created_at": {
                                      "gte": "%s.000+0800",
                                      "lte": "%s.000+0800"
                                    }
                                  }
                                },
                                {
                                  "query_string": {
                                    "analyze_wildcard": true,
                                    "query": "*"
                                  }
                                }
                              ]
                            }
                          },
                          "aggs": {
                            "2": {
                              "terms": {
                                "field": "html_url.keyword",
                                "size": 500,
                                "order": {
                                  "_key": "desc"
                                },
                                "min_doc_count": 1
                              },
                              "aggs": {
                                "3": {
                                  "max": {
                                    "field": "stargazers_count"
                                  }
                                }
                              }
                            }
                          }
                        }''' % (starttime, endtime)
            url = self.github_star
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["github_star"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      }
                    }
                  },"aggs": {
                    "user_login": {
                      "sum": {
                        "field": "is_gitee_star"}}}
                }''' % (starttime, endtime)
            url = self.gitee_star_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["gitee_star_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          }
                        }
                      },"aggs": {
                        "user_login": {
                          "sum": {
                            "field": "is_gitee_fork"}}}
                    }''' % (starttime, endtime)
            url = self.gitee_fork_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["gitee_fork_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          }
                        }
                      },"aggs": {
                        "user_login": {
                          "sum": {
                            "field": "is_gitee_watch"}}}
                    }''' % (starttime, endtime)
            url = self.gitee_watch_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["gitee_watch_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                          "query": {
                            "bool": {
                              "filter":{
                                "range":{
                                  "created_at":{
                                    "gte":"%s.000+0800",
                                    "lt":"%s.000+0800"
                                  }
                                }
                              }
                            }
                          },"aggs": {
                            "user_login": {
                              "sum": {
                                "field": "is_gitee_pull_request"}}}
                        }''' % (starttime, endtime)
            url = self.gitee_pr_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["gitee_pr_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                              "query": {
                                "bool": {
                                  "filter":{
                                    "range":{
                                      "created_at":{
                                        "gte":"%s.000+0800",
                                        "lt":"%s.000+0800"
                                      }
                                    }
                                  }
                                }
                              },"aggs": {
                                "user_login": {
                                  "sum": {
                                    "field": "is_gitee_issue"}}}
                            }''' % (starttime, endtime)
            url = self.gitee_issue
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["gitee_issue"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{
              "size": 0,
              "query": {
                "bool": {
                  "filter": [
                    {
                      "range": {
                        "created_at": {
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                        }
                      }
                    },
                    {
                      "query_string": {
                        "analyze_wildcard": true,
                        "query": "(path:\\"basicdb.*\\" OR path:\\"opengauss.*\\") AND !path:\\"odbc.tar.gz.sha256\\" AND !path:\\"64bit.tar.gz.sha256\\"  AND !path:\\"jdbc.tar.gz.sha256\\" AND !path:\\"symbol.tar.gz.sha256\\" AND !path:\\"libpq.tar.gz.sha256\\" AND !path.keyword:\\"1.0.0/openGauss-third_party_binarylibs.tar.gz/info/refs\\" AND !path.keyword:\\"1.0.0/openGauss-third\\""
                      }
                    }
                  ]
                }
              },
              "aggs": {
                "2": {
                  "date_histogram": {
                    "interval": "1000d",
                    "field": "created_at",
                    "min_doc_count": 0,
                    "extended_bounds": {
                      "min": 1590940800000,
                      "max": 1604137994078
                    },
                    "format": "epoch_millis"
                  },
                  "aggs": {
                    "1": {
                      "cardinality": {
                        "field": "ip.keyword"
                      }
                    }
                  }
                }
              }
            }''' % (starttime, endtime)
            url = self.D0_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["D0_user"] = int(res["aggregations"]["2"]["buckets"][0]['1']['value'])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },
                    // "must": [{ "match": { "ip":"127.0.0.1" }}],
                    "should": [{
                            "match_phrase": {
                                "is_gitee_fork": 1
                            }
                        },
                  {
                            "match_phrase": {
                                "is_gitee_issue": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_pull_request": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_review_comment": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_comment": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_issue_comment": 1
                            }
                        }],
                        "minimum_should_match": 1
                }
              },"aggs": {
                "user_login": {
                  "cardinality": {
                    "field": "user_login.keyword"}}}
            }''' % (starttime, endtime)
            url = self.contributor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["contributor"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      },
                        // "must": [{ "match": { "ip":"127.0.0.1" }}],
                        "should": [
                      {
                                "match_phrase": {
                                    "is_gitee_issue": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_pull_request": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_review_comment": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_comment": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_issue_comment": 1
                                }
                            }],
                            "minimum_should_match": 1
                    }
                  },"aggs": {
                    "user_login": {
                      "cardinality": {
                        "field": "user_login.keyword"}}}
                }''' % (starttime, endtime)
            url = self.deep_contributor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["deep_contributor"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },
                    "must": [{ "match": { "is_project_internal_user":1 }}],
                    "should": [{
                            "match_phrase": {
                                "is_gitee_fork": 1
                            }
                        },
                  {
                            "match_phrase": {
                                "is_gitee_issue": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_pull_request": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_review_comment": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_comment": 1
                            }
                        },
                        {
                            "match_phrase": {
                                "is_gitee_issue_comment": 1
                            }
                        }],
                        "minimum_should_match": 1
                }
              },"aggs": {
                "user_login": {
                  "cardinality": {
                    "field": "user_login.keyword"}}}
            }''' % (starttime, endtime)
            url = self.contributorHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["contributorHuawei"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      },
                        "must": [{ "match": { "is_project_internal_user":0 }}],
                        "should": [{
                                "match_phrase": {
                                    "is_gitee_fork": 1
                                }
                            },
                      {
                                "match_phrase": {
                                    "is_gitee_issue": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_pull_request": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_review_comment": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_comment": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_issue_comment": 1
                                }
                            }],
                            "minimum_should_match": 1
                    }
                  },"aggs": {
                    "user_login": {
                      "cardinality": {
                        "field": "user_login.keyword"}}}
                }''' % (starttime, endtime)
            url = self.contributorNotHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["contributorNotHuawei"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },"aggs": {
                "user_login": {
                  "cardinality": {
                    "field": "user_gitee_name.keyword"}}}
            }''' % (starttime, endtime)
            url = self.committers
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["committers"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      },"must": [{ "match": { "is_project_internal_user":1 }}]
                    }
                  },"aggs": {
                    "user_login": {
                      "cardinality": {
                        "field": "user_gitee_name.keyword"}}}
                }''' % (starttime, endtime)
            url = self.committersHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["committersHuawei"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      },"must_not": [{ "match": { "is_project_internal_user":1 }}]
                    }
                  },"aggs": {
                    "user_gitee_name": {
                      "cardinality": {
                        "field": "user_gitee_name.keyword"}}}
                }''' % (starttime, endtime)
            url = self.committersNotHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["opengauss"]["committersNotHuawei"] = int(res["aggregations"]["user_gitee_name"]["value"])

        elif org == 'openlookeng':
            # openlookeng
            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  },
                    // "must": [{ "match": { "ip":"127.0.0.1" }}],
                    "should": [{
                            "match_phrase": {
                                "source_type_title": "搜索引擎"
                            }
                        },
                  {
                            "match_phrase": {
                                "source_type_title": "外部链接"
                            }
                        },
                        {
                            "match_phrase": {
                                "source_type_title": "直接访问"
                            }
                        }],
                        "minimum_should_match": 1
                }
              },"aggs": {
                "ip_count": {
                  "sum": {
                    "field": "ip_count"}}}
            }''' % (starttime, endtime)
            url = self.webvisitor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["webvisitor"] = result["openlookeng"]["D0_user"] = int(res["aggregations"]["ip_count"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      },
                        // "must": [{ "match": { "ip":"127.0.0.1" }}],
                        "should": [{
                                "match_phrase": {
                                    "source_type_title": "搜索引擎"
                                }
                            },
                      {
                                "match_phrase": {
                                    "source_type_title": "外部链接"
                                }
                            },
                            {
                                "match_phrase": {
                                    "source_type_title": "直接访问"
                                }
                            }],
                            "minimum_should_match": 1
                    }
                  },"aggs": {
                    "pv_count": {
                      "sum": {
                        "field": "pv_count"}}}
                }''' % (starttime, endtime)
            url = self.website
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["website"] = int(res["aggregations"]["pv_count"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },"aggs": {
                "user_id": {
                  "cardinality": {
                    "field": "user_id.keyword"}}}
            }''' % (starttime, endtime)
            url = self.maillist
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["maillist"] = int(res["aggregations"]["user_id"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "metadata__updated_on":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      }, "must": [{ "match": {"project":"openlookeng_mail" }}]
                    }
                  },"aggs": {
                    "user_id": {
                      "sum": {
                        "field": "is_hyperkitty_message"}}}
                }''' % (starttime, endtime)
            url = self.maillistnum
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["maillistnum"] = int(res["aggregations"]["user_id"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }, "must": [{ "match": { "event":"DOWNLOAD ZIP" }},{ "match": { "is_forked_repo":1 }}]
                }
              },
              "aggs": {
                "author_name": {
                  "cardinality": {
                    "field": "author_name.keyword"}}}
            }''' % (starttime, endtime)
            url = self.repo
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["repo"] = int(res["aggregations"]["author_name"]["value"])

            data = '''{"size":10000,
              "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }, "must_not": [{ "match": { "is_topic":1 }}]
                }
              },
              "aggs": {
                "new_followers": {
                  "sum": {
                    "field": "new_followers"}}}
            }''' % (starttime, endtime)
            url = self.video
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["video"] = int(res["aggregations"]["new_followers"]["value"])

            data = '''{"size":100,
             "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },
                  "aggs": {
                    "user_login": {
                      "sum": {
                        "field": "is_gitee_star"}}}
                }'''% (starttime, endtime)
            url = self.gitee_star_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["gitee_star_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":100,
            "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },
                     "aggs": {
                        "user_login": {
                          "sum": {
                            "field": "is_gitee_fork"}}}
                    }'''% (starttime, endtime)
            url = self.gitee_fork_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["gitee_fork_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
            "query": {
                "bool": {
                  "filter":{
                    "range":{
                      "created_at":{
                        "gte":"%s.000+0800",
                        "lt":"%s.000+0800"
                      }
                    }
                  }
                }
              },
                     "aggs": {
                        "user_login": {
                          "sum": {
                            "field": "is_gitee_watch"}}}
                    }'''% (starttime, endtime)
            url = self.gitee_watch_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["gitee_watch_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                          "query": {
                            "bool": {
                              "filter":{
                                "range":{
                                  "created_at":{
                                    "gte":"%s.000+0800",
                                    "lt":"%s.000+0800"
                                  }
                                }
                              }
                            }
                          },"aggs": {
                            "user_login": {
                              "sum": {
                                "field": "is_gitee_pull_request"}}}
                        }''' % (starttime, endtime)
            url = self.gitee_pr_user
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["gitee_pr_user"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{
                  "size": 0,
                  "query": {
                    "bool": {
                      "filter": [
                        {
                          "range": {
                            "created_at": {
                              "gte": "%s.000+0800",
                              "lte": "%s.000+0800"
                            }
                          }
                        },
                        {
                          "query_string": {
                            "analyze_wildcard": true,
                            "query": "*"
                          }
                        }
                      ]
                    }
                  },
                  "aggs": {
                    "2": {
                      "terms": {
                        "field": "html_url.keyword",
                        "size": 500,
                        "order": {
                          "_key": "desc"
                        },
                        "min_doc_count": 1
                      },
                      "aggs": {
                        "3": {
                          "max": {
                            "field": "forks_count"
                          }
                        }
                      }
                    }
                  }
                }''' % (starttime, endtime)
            url = self.github_fork
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["github_fork"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{
                      "size": 0,
                      "query": {
                        "bool": {
                          "filter": [
                            {
                              "range": {
                                "created_at": {
                                  "gte": "%s.000+0800",
                                  "lte": "%s.000+0800"
                                }
                              }
                            },
                            {
                              "query_string": {
                                "analyze_wildcard": true,
                                "query": "*"
                              }
                            }
                          ]
                        }
                      },
                      "aggs": {
                        "2": {
                          "terms": {
                            "field": "html_url.keyword",
                            "size": 500,
                            "order": {
                              "_key": "desc"
                            },
                            "min_doc_count": 1
                          },
                          "aggs": {
                            "3": {
                              "max": {
                                "field": "subscribers_count"
                              }
                            }
                          }
                        }
                      }
                    }''' % (starttime, endtime)
            url = self.github_watch
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["github_watch"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{
                          "size": 0,
                          "query": {
                            "bool": {
                              "filter": [
                                {
                                  "range": {
                                    "created_at": {
                                      "gte": "%s.000+0800",
                                      "lte": "%s.000+0800"
                                    }
                                  }
                                },
                                {
                                  "query_string": {
                                    "analyze_wildcard": true,
                                    "query": "*"
                                  }
                                }
                              ]
                            }
                          },
                          "aggs": {
                            "2": {
                              "terms": {
                                "field": "html_url.keyword",
                                "size": 500,
                                "order": {
                                  "_key": "desc"
                                },
                                "min_doc_count": 1
                              },
                              "aggs": {
                                "3": {
                                  "max": {
                                    "field": "stargazers_count"
                                  }
                                }
                              }
                            }
                          }
                        }''' % (starttime, endtime)
            url = self.github_star
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["github_star"] = sum([int(r['3']['value']) for r in res["aggregations"]["2"]["buckets"]])

            data = '''{"size":10000,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          }
                        }
                      },"aggs": {
                        "user_login": {
                          "sum": {
                            "field": "is_gitee_issue"}}}
                    }''' % (starttime, endtime)
            url = self.gitee_issue
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["gitee_issue"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      },
                        // "must": [{ "match": { "ip":"127.0.0.1" }}],
                        "should": [{
                                "match_phrase": {
                                    "is_gitee_fork": 1
                                }
                            },
                      {
                                "match_phrase": {
                                    "is_gitee_issue": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_pull_request": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_review_comment": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_comment": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_issue_comment": 1
                                }
                            }],
                            "minimum_should_match": 1
                    }
                  },"aggs": {
                    "user_login": {
                      "cardinality": {
                        "field": "user_login.keyword"}}}
                }''' % (starttime, endtime)
            url = self.contributor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["contributor"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          },
                            // "must": [{ "match": { "ip":"127.0.0.1" }}],
                            "should": [
                          {
                                    "match_phrase": {
                                        "is_gitee_issue": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_pull_request": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_review_comment": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_comment": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_issue_comment": 1
                                    }
                                }],
                                "minimum_should_match": 1
                        }
                      },"aggs": {
                        "user_login": {
                          "cardinality": {
                            "field": "user_login.keyword"}}}
                    }''' % (starttime, endtime)
            url = self.deep_contributor
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["deep_contributor"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      },
                        "must": [{ "match": { "is_project_internal_user":1 }}],
                        "should": [{
                                "match_phrase": {
                                    "is_gitee_fork": 1
                                }
                            },
                      {
                                "match_phrase": {
                                    "is_gitee_issue": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_pull_request": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_review_comment": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_comment": 1
                                }
                            },
                            {
                                "match_phrase": {
                                    "is_gitee_issue_comment": 1
                                }
                            }],
                            "minimum_should_match": 1
                    }
                  },"aggs": {
                    "user_login": {
                      "cardinality": {
                        "field": "user_login.keyword"}}}
                }''' % (starttime, endtime)
            url = self.contributorHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["contributorHuawei"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          },
                            "must": [{ "match": { "is_project_internal_user":0 }}],
                            "should": [{
                                    "match_phrase": {
                                        "is_gitee_fork": 1
                                    }
                                },
                          {
                                    "match_phrase": {
                                        "is_gitee_issue": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_pull_request": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_review_comment": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_comment": 1
                                    }
                                },
                                {
                                    "match_phrase": {
                                        "is_gitee_issue_comment": 1
                                    }
                                }],
                                "minimum_should_match": 1
                        }
                      },"aggs": {
                        "user_login": {
                          "cardinality": {
                            "field": "user_login.keyword"}}}
                    }''' % (starttime, endtime)
            url = self.contributorNotHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["contributorNotHuawei"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                  "query": {
                    "bool": {
                      "filter":{
                        "range":{
                          "created_at":{
                            "gte":"%s.000+0800",
                            "lt":"%s.000+0800"
                          }
                        }
                      }
                    }
                  },"aggs": {
                    "user_login": {
                      "cardinality": {
                        "field": "user_gitee_name.keyword"}}}
                }''' % (starttime, endtime)
            url = self.committers
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["committers"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          },"must_not": [{ "match": { "is_project_internal_user":1 }}]
                        }
                      },"aggs": {
                        "user_login": {
                          "cardinality": {
                            "field": "user_gitee_name.keyword"}}}
                    }''' % (starttime, endtime)
            url = self.contributorHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["committersHuawei"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{"size":10000,
                      "query": {
                        "bool": {
                          "filter":{
                            "range":{
                              "created_at":{
                                "gte":"%s.000+0800",
                                "lt":"%s.000+0800"
                              }
                            }
                          },"must": [{ "match": { "is_project_internal_user":1 }}]
                        }
                      },"aggs": {
                        "user_login": {
                          "cardinality": {
                            "field": "user_gitee_name.keyword"}}}
                    }''' % (starttime, endtime)
            url = self.contributorNotHuawei
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["committersNotHuawei"] = int(res["aggregations"]["user_login"]["value"])

            data = '''{
            "query": {
            "bool": {
                    "must_not": [{ "match": { "path.keyword":"-" }}]
                }
                }
        }'''
            url = self.paly_video
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["paly_video"] = int(res["count"])

            data = '''{
                        "query": {
                        "bool": {
                        "filter":{
                            "range":{
                              "created_at": {
                                  "gte": "%s.000+0800",
                                  "lte": "%s.000+0800"
                              }
                            }
                          },
                                "must": [{ "match": { "hostname":"download.openlookeng.io" }}]
                            }
                            }
                    }'''% (starttime, endtime)
            url = self.downlaod
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["downlaod"] = int(res["count"])
            result["openlookeng"]["downlaod_gb"] = "N/A"

            data = '''{"size":100,
                            "query": {
                                "bool": {
                                  "filter":{
                                    "range":{
                                      "created_at":{
                                        "gte":"%s.000+0800",
                                        "lt":"%s.000+0800"
                                      }
                                    }
                                  },"must": [{ "match": { "hostname":"download.openlookeng.io" }}]
                                }
                              },
                                     "aggs": {
                                        "ip": {
                                          "cardinality": {
                                            "field": "ip.keyword"}}}
                                    }''' % (starttime, endtime)
            url = self.downlaod_ip
            res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
            result["openlookeng"]["downlaod_ip"] = int(res["aggregations"]["ip"]["value"])

        return result

    def writePPT(self, data, itime, org=None, fileppt='', data2={}):
        def _addText(prs, content, slides, inches, size, color=''):
            textbox = prs.slides[slides].shapes.add_textbox(Inches(inches[0]), Inches(inches[1]), Inches(inches[2]),
                                                            Inches(inches[3]))
            tf = textbox.text_frame
            para = tf.add_paragraph()
            para.text = str(content)
            font = para.font
            font.size = Pt(size)
            if color == 'red':
                font.color.rgb = RGBColor(255, 0, 0)
            elif color == 'green':
                font.color.rgb = RGBColor(0, 255, 0)

        prs = Presentation('D:\\模板.pptx')

        # add picture
        # openeuler
        if not org or org == 'openeuler':
            prs.slides[3].shapes.add_picture('D:\\GrafanaPng\\openeuler_webviews.png', Inches(1), Inches(1.5), Inches(11.3),
                                             Inches(3))
            prs.slides[3].shapes.add_picture('D:\\GrafanaPng\\openeuler_maillist.png', Inches(1), Inches(4.72), Inches(6),
                                             Inches(2))
            prs.slides[3].shapes.add_picture('D:\\GrafanaPng\\openeuler_repo.png', Inches(7.1), Inches(4.72), Inches(6),
                                             Inches(2))
            prs.slides[4].shapes.add_picture('D:\\GrafanaPng\\openeuler_users.png', Inches(1), Inches(1.5), Inches(11.3),
                                             Inches(5))
            prs.slides[5].shapes.add_picture('D:\\GrafanaPng\\openeuler_contributorlines.png', Inches(1), Inches(1.5),
                                             Inches(11.3), Inches(5))
            prs.slides[6].shapes.add_picture('D:\\GrafanaPng\\openeuler_committerlines.png', Inches(1), Inches(1.5),
                                             Inches(11.3), Inches(5))

        # mindspore
        if not org or org == 'mindspore':
            prs.slides[9].shapes.add_picture('D:\\GrafanaPng\\mindspore_webviews.png', Inches(1), Inches(1.5), Inches(11.3),
                                             Inches(3))
            prs.slides[9].shapes.add_picture('D:\\GrafanaPng\\mindspore_maillist.png', Inches(1), Inches(4.72), Inches(6),
                                             Inches(2))
            prs.slides[9].shapes.add_picture('D:\\GrafanaPng\\mindspore_repo.png', Inches(7.1), Inches(4.72), Inches(6),
                                             Inches(2))
            prs.slides[10].shapes.add_picture('D:\\GrafanaPng\\mindspore_users.png', Inches(1), Inches(1.5), Inches(11.3),
                                             Inches(5))
            prs.slides[11].shapes.add_picture('D:\\GrafanaPng\\mindspore_contributorlines.png', Inches(1), Inches(1.5),
                                             Inches(11.3), Inches(5))
            prs.slides[12].shapes.add_picture('D:\\GrafanaPng\\mindspore_committerlines.png', Inches(1), Inches(1.5),
                                              Inches(11.3), Inches(5))

        # opengauss
        if not org or org == 'opengauss':
            prs.slides[15].shapes.add_picture('D:\\GrafanaPng\\opengauss_webviews.png', Inches(1), Inches(1.5), Inches(11.3),
                                              Inches(3))
            prs.slides[15].shapes.add_picture('D:\\GrafanaPng\\opengauss_maillist.png', Inches(1), Inches(4.72), Inches(6),
                                              Inches(2))
            prs.slides[15].shapes.add_picture('D:\\GrafanaPng\\opengauss_repo.png', Inches(7.1), Inches(4.72), Inches(6),
                                              Inches(2))
            prs.slides[16].shapes.add_picture('D:\\GrafanaPng\\opengauss_users.png', Inches(1), Inches(1.5), Inches(11.3),
                                              Inches(5))
            prs.slides[17].shapes.add_picture('D:\\GrafanaPng\\opengauss_contributorlines.png', Inches(1), Inches(1.5),
                                              Inches(11.3), Inches(5))
            prs.slides[18].shapes.add_picture('D:\\GrafanaPng\\opengauss_committerlines.png', Inches(1), Inches(1.5),
                                              Inches(11.3), Inches(5))

        # openlookeng
        if not org or org == 'openlookeng':
            prs.slides[21].shapes.add_picture('D:\\GrafanaPng\\openlookeng_webviews.png', Inches(1), Inches(1.5), Inches(11.3),
                                              Inches(3))
            prs.slides[21].shapes.add_picture('D:\\GrafanaPng\\openlookeng_maillist.png', Inches(1), Inches(4.72), Inches(6),
                                              Inches(2))
            prs.slides[21].shapes.add_picture('D:\\GrafanaPng\\openlookeng_repo.png', Inches(7.1), Inches(4.72), Inches(6),
                                              Inches(2))
            prs.slides[22].shapes.add_picture('D:\\GrafanaPng\\openlookeng_users.png', Inches(1), Inches(1.5), Inches(11.3),
                                              Inches(5))
            prs.slides[23].shapes.add_picture('D:\\GrafanaPng\\openlookeng_contributorlines.png', Inches(1), Inches(1.5),
                                              Inches(11.3), Inches(5))
            prs.slides[24].shapes.add_picture('D:\\GrafanaPng\\openlookeng_committerlines.png', Inches(1), Inches(1.5),
                                              Inches(11.3), Inches(5))

        # add text
        # openeuler
        home_index_openeuler = 1
        compare_index_openeuler = 2
        if not org or org == 'openeuler':
            _addText(prs, itime, home_index_openeuler, [3.397, 0.476, 1, 1], 14)
            _addText(prs, data['openeuler']['website'], home_index_openeuler, [5.582, 2.35, 1, 1], 16)  # 触达用户
            _addText(prs, data['openeuler']['website'], home_index_openeuler, [7.114, 2.232, 1, 1], 14)  # 网站访问
            _addText(prs, data["openeuler"]["maillist"], home_index_openeuler, [8.421, 2.338, 1, 1], 14)  # 邮件列表
            _addText(prs, data["openeuler"]["repo"], home_index_openeuler, [9.37, 2.338, 1, 1], 14)  # 代码仓
            _addText(prs, data["openeuler"]["video"], home_index_openeuler, [10.472, 2.338, 1, 1], 14)  # 视频
            _addText(prs, data["openeuler"]["downloadall"], home_index_openeuler, [4.968, 3.437, 1, 1], 16)  # user
            _addText(prs, data["openeuler"]["downloadall"], home_index_openeuler, [7.236, 3.496, 1, 1], 14)  # 试用
            _addText(prs, data["openeuler"]["contributor"], home_index_openeuler, [4.645, 4.527, 1, 1], 16)  # Contributor
            _addText(prs, data["openeuler"]["contributorHuawei"], home_index_openeuler, [7.318, 4.704, 1, 1], 14)  # 华为
            _addText(prs, data["openeuler"]["contributorNotHuawei"], home_index_openeuler, [9.641, 4.633, 1, 1], 14)  # 非华为
            _addText(prs, data["openeuler"]["committers"], home_index_openeuler, [4.181, 5.732, 1, 1], 16)  # Committer
            _addText(prs, data["openeuler"]["committersHuawei"], home_index_openeuler, [7.377, 5.984, 1, 1], 14)  # 华为
            _addText(prs, data["openeuler"]["committersNotHuawei"], home_index_openeuler, [9.685, 6.055, 1, 1], 14)  # 非华为
            # compare
            _addText(prs, data['openeuler']['webvisitor'], compare_index_openeuler, [2.795, 1.098, 1, 1], 10, 'green')
            _addText(prs, data['openeuler']['website'], compare_index_openeuler, [2.795, 1.48, 1, 1], 10, 'green')
            _addText(prs, 0, compare_index_openeuler, [2.795, 1.834, 1, 1], 10, 'green')
            _addText(prs, data["openeuler"]["maillist"], compare_index_openeuler, [2.795, 2.208, 1, 1], 10, 'green')
            _addText(prs, data["openeuler"]["maillistnum"], compare_index_openeuler, [2.795, 2.582, 1, 1], 10, 'green')
            _addText(prs, 0, compare_index_openeuler, [2.795, 2.975, 1, 1], 10, 'green')
            _addText(prs, data2['openeuler']['webvisitor'], compare_index_openeuler, [5.263, 1.098, 1, 1], 10, 'red')
            _addText(prs, data2['openeuler']['website'], compare_index_openeuler, [5.263, 1.48, 1, 1], 10, 'red')
            _addText(prs, 0, compare_index_openeuler, [5.263, 1.834, 1, 1], 10, 'red')
            _addText(prs, data2['openeuler']['maillist'], compare_index_openeuler, [5.263, 2.208, 1, 1], 10, 'red')
            _addText(prs, data2['openeuler']['maillistnum'], compare_index_openeuler, [5.263, 2.582, 1, 1], 10, 'red')
            _addText(prs, 0, compare_index_openeuler, [5.263, 2.975, 1, 1], 10, 'red')
            _addText(prs, 0, compare_index_openeuler, [8.356, 1.968, 1, 1], 10, 'green')
            _addText(prs, data["openeuler"]["video"], compare_index_openeuler, [8.356, 1.602, 1, 1], 10, 'green')
            _addText(prs, 0, compare_index_openeuler, [8.356, 1.251, 1, 1], 10, 'green')
            _addText(prs, 0, compare_index_openeuler, [10.338, 1.968, 1, 1], 10, 'red')
            _addText(prs, data2["openeuler"]["video"], compare_index_openeuler, [10.338, 1.602, 1, 1], 10, 'red')
            _addText(prs, 0, compare_index_openeuler, [10.338, 1.251, 1, 1], 10, 'red')
            _addText(prs, data['openeuler']['website'], compare_index_openeuler, [2.877, 3.626, 1, 1], 10, 'green')
            _addText(prs, data["openeuler"]["contributor"], compare_index_openeuler, [2.877, 3.992, 1, 1], 10, 'green')
            _addText(prs, data["openeuler"]["deep_contributor"], compare_index_openeuler, [2.877, 4.358, 1, 1], 10, 'green')
            _addText(prs, data["openeuler"]["gitee_star_user"]+data["openeuler"]['github_star'], compare_index_openeuler, [2.877, 4.732, 1, 1], 10, 'green')
            _addText(prs, data["openeuler"]["gitee_fork_user"]+data["openeuler"]['github_fork'], compare_index_openeuler, [2.877, 5.106, 1, 1], 10, 'green')
            _addText(prs, data["openeuler"]["gitee_watch_user"]+data["openeuler"]['github_watch'], compare_index_openeuler, [2.877, 5.472, 1, 1], 10, 'green')
            _addText(prs, data["openeuler"]["gitee_star_user"], compare_index_openeuler, [2.877, 5.838, 1, 1], 10, 'green')
            _addText(prs, data["openeuler"]["gitee_pr_user"], compare_index_openeuler, [2.877, 6.173, 1, 1], 10, 'green')
            _addText(prs, data["openeuler"]["gitee_issue"], compare_index_openeuler, [2.877, 6.501, 1, 1], 10, 'green')
            _addText(prs, data2['openeuler']['website'], compare_index_openeuler, [5.251, 3.626, 1, 1], 10, 'red')
            _addText(prs, data2["openeuler"]["contributor"], compare_index_openeuler, [5.251, 3.992, 1, 1], 10, 'red')
            _addText(prs, data2["openeuler"]["committers"], compare_index_openeuler, [5.251, 4.358, 1, 1], 10, 'red')
            _addText(prs, data2["openeuler"]["gitee_star_user"]+data2["openeuler"]['github_star'], compare_index_openeuler, [5.251, 4.732, 1, 1], 10, 'red')
            _addText(prs, data2["openeuler"]["gitee_fork_user"]+data2["openeuler"]['github_fork'], compare_index_openeuler, [5.251, 5.106, 1, 1], 10, 'red')
            _addText(prs, data2["openeuler"]["gitee_watch_user"]+data2["openeuler"]['github_watch'], compare_index_openeuler, [5.251, 5.472, 1, 1], 10, 'red')
            _addText(prs, data2["openeuler"]["gitee_star_user"], compare_index_openeuler, [5.251, 5.838, 1, 1], 10, 'red')
            _addText(prs, data2["openeuler"]["gitee_pr_user"], compare_index_openeuler, [5.251, 6.173, 1, 1], 10, 'red')
            _addText(prs, data2["openeuler"]["gitee_issue"], compare_index_openeuler, [5.251, 6.501, 1, 1], 10, 'red')
        # mindspore
        home_index_mindspore = 7
        compare_index_mindspore = 8
        if not org or org == 'mindspore':
            _addText(prs, itime, home_index_mindspore, [3.397, 0.476, 1, 1], 14)
            _addText(prs, data['mindspore']['website'], home_index_mindspore, [5.582, 2.35, 1, 1], 16)  # 触达用户
            _addText(prs, data['mindspore']['website'], home_index_mindspore, [7.114, 2.232, 1, 1], 14)  # 网站访问
            _addText(prs, data["mindspore"]["maillist"], home_index_mindspore, [8.421, 2.338, 1, 1], 14)  # 邮件列表
            _addText(prs, data["mindspore"]["repo"], home_index_mindspore, [9.37, 2.338, 1, 1], 14)  # 代码仓
            _addText(prs, data["mindspore"]["video"], home_index_mindspore, [10.472, 2.338, 1, 1], 14)  # 视频
            _addText(prs, data["mindspore"]["users"], home_index_mindspore, [4.968, 3.437, 1, 1], 16)  # user
            _addText(prs, data["mindspore"]["users"], home_index_mindspore, [7.236, 3.496, 1, 1], 14)  # 试用
            _addText(prs, data["mindspore"]["contributor"], home_index_mindspore, [4.645, 4.527, 1, 1], 16)  # Contributor
            _addText(prs, data["mindspore"]["contributorHuawei"], home_index_mindspore, [7.318, 4.704, 1, 1], 14)  # 华为
            _addText(prs, data["mindspore"]["contributorNotHuawei"], home_index_mindspore, [9.641, 4.633, 1, 1], 14)  # 非华为
            _addText(prs, data["mindspore"]["committers"], home_index_mindspore, [4.181, 5.732, 1, 1], 16)  # Committer
            _addText(prs, data["mindspore"]["committers"], home_index_mindspore, [7.377, 5.984, 1, 1], 14)  # 华为
            _addText(prs, '0', home_index_mindspore, [9.685, 6.055, 1, 1], 14)  # 非华为
            # compare
            _addText(prs, data["mindspore"]['webvisitor'], compare_index_mindspore, [2.795, 1.098, 1, 1], 10, 'green')
            _addText(prs, data["mindspore"]['website'], compare_index_mindspore, [2.795, 1.48, 1, 1], 10, 'green')
            _addText(prs, 0, compare_index_mindspore, [2.795, 1.834, 1, 1], 10, 'green')
            _addText(prs, data["mindspore"]["maillist"], compare_index_mindspore, [2.795, 2.208, 1, 1], 10, 'green')
            _addText(prs, data["mindspore"]["maillistnum"], compare_index_mindspore, [2.795, 2.582, 1, 1], 10, 'green')
            _addText(prs, data2["mindspore"]['webvisitor'], compare_index_mindspore, [5.263, 1.098, 1, 1], 10, 'red')
            _addText(prs, data2["mindspore"]['website'], compare_index_mindspore, [5.263, 1.48, 1, 1], 10, 'red')
            _addText(prs, data2["mindspore"]['repo'], compare_index_mindspore, [5.263, 1.834, 1, 1], 10, 'red')
            _addText(prs, data2["mindspore"]['maillist'], compare_index_mindspore, [5.263, 2.208, 1, 1], 10, 'red')
            _addText(prs, data2["mindspore"]['maillistnum'], compare_index_mindspore, [5.263, 2.582, 1, 1], 10, 'red')
            _addText(prs, '0', 21, [8.356, 1.968, 1, 1], 10, 'green')
            _addText(prs, data["mindspore"]["video"], compare_index_mindspore, [8.356, 1.602, 1, 1], 10, 'green')
            _addText(prs, '0', 21, [8.356, 1.251, 1, 1], 10, 'green')
            _addText(prs, '0', 21, [10.338, 1.968, 1, 1], 10, 'red')
            _addText(prs, data2["mindspore"]["video"], compare_index_mindspore, [10.338, 1.602, 1, 1], 10, 'red')
            _addText(prs, '0', 21, [10.338, 1.251, 1, 1], 10, 'red')
            _addText(prs, data["mindspore"]["contributor"], compare_index_mindspore, [2.877, 3.992, 1, 1], 10, 'green')
            _addText(prs, data["mindspore"]["committers"], compare_index_mindspore, [2.877, 4.358, 1, 1], 10, 'green')
            _addText(prs, data["mindspore"]["gitee_star_user"], compare_index_mindspore, [2.877, 4.732, 1, 1], 10, 'green')
            _addText(prs, data["mindspore"]["gitee_fork_user"], compare_index_mindspore, [2.877, 5.106, 1, 1], 10, 'green')
            _addText(prs, data["mindspore"]["gitee_watch_user"], compare_index_mindspore, [2.877, 5.472, 1, 1], 10, 'green')
            _addText(prs, data["mindspore"]["gitee_star_user"], compare_index_mindspore, [2.877, 5.838, 1, 1], 10, 'green')
            _addText(prs, data["mindspore"]["gitee_pr_user"], compare_index_mindspore, [2.877, 6.173, 1, 1], 10, 'green')
            _addText(prs, data2["mindspore"]["contributor"], compare_index_mindspore, [5.251, 3.992, 1, 1], 10, 'red')
            _addText(prs, data2["mindspore"]["committers"], compare_index_mindspore, [5.251, 4.358, 1, 1], 10, 'red')
            _addText(prs, data2["mindspore"]["gitee_star_user"], compare_index_mindspore, [5.251, 4.732, 1, 1], 10, 'red')
            _addText(prs, data2["mindspore"]["gitee_fork_user"], compare_index_mindspore, [5.251, 5.106, 1, 1], 10, 'red')
            _addText(prs, data2["mindspore"]["gitee_watch_user"], compare_index_mindspore, [5.251, 5.472, 1, 1], 10, 'red')
            _addText(prs, data2["mindspore"]["gitee_star_user"], compare_index_mindspore, [5.251, 5.838, 1, 1], 10, 'red')
            _addText(prs, data2["mindspore"]["gitee_pr_user"], compare_index_mindspore, [5.251, 6.173, 1, 1], 10, 'red')
        # opengauss
        home_index_opengauss = 13
        compare_index_opengauss = 14
        if not org or org == 'opengauss':
            _addText(prs, itime, home_index_opengauss, [3.397, 0.476, 1, 1], 14)  # 触达用户
            _addText(prs, data['opengauss']['website'], home_index_opengauss, [5.582, 2.35, 1, 1], 16)  # 触达用户
            _addText(prs, data['opengauss']['website'], home_index_opengauss, [7.114, 2.232, 1, 1], 14)  # 网站访问
            _addText(prs, data["opengauss"]["maillist"], home_index_opengauss, [8.421, 2.338, 1, 1], 14)  # 邮件列表
            _addText(prs, data["opengauss"]["repo"], home_index_opengauss, [9.37, 2.338, 1, 1], 14)  # 代码仓
            _addText(prs, data["opengauss"]["video"], home_index_opengauss, [10.472, 2.338, 1, 1], 14)  # 视频
            _addText(prs, data["opengauss"]["users"], home_index_opengauss, [4.968, 3.437, 1, 1], 16)  # user
            _addText(prs, data["opengauss"]["users"], home_index_opengauss, [7.236, 3.496, 1, 1], 14)  # 试用
            _addText(prs, data["opengauss"]["contributor"], home_index_opengauss, [4.645, 4.527, 1, 1], 16)  # Contributor
            _addText(prs, data["opengauss"]["contributorHuawei"], home_index_opengauss, [7.318, 4.704, 1, 1], 14)  # 华为
            _addText(prs, data["opengauss"]["contributorNotHuawei"], home_index_opengauss, [9.641, 4.633, 1, 1], 14)  # 非华为
            _addText(prs, data["opengauss"]["committers"], home_index_opengauss, [4.181, 5.732, 1, 1], 16)  # Committer
            _addText(prs, data["opengauss"]["committersHuawei"], home_index_opengauss, [7.377, 5.984, 1, 1], 14)  # 华为
            _addText(prs, data["opengauss"]["committersNotHuawei"], home_index_opengauss, [9.685, 6.055, 1, 1], 14)  # 非华为
            # compare
            _addText(prs, data["opengauss"]['webvisitor'], compare_index_opengauss, [2.795, 1.098, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]['website'], compare_index_opengauss, [2.795, 1.48, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["downlaod"], compare_index_opengauss, [2.795, 1.834, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["maillist"], compare_index_opengauss, [2.795, 2.208, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["maillistnum"], compare_index_opengauss, [2.795, 2.582, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["paly_video"], compare_index_opengauss, [2.795, 2.975, 1, 1], 10, 'green')
            _addText(prs, data2["opengauss"]['webvisitor'], compare_index_opengauss, [5.263, 1.098, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]['website'], compare_index_opengauss, [5.263, 1.48, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]['downlaod'], compare_index_opengauss, [5.263, 1.834, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]['maillist'], compare_index_opengauss, [5.263, 2.208, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]['maillistnum'], compare_index_opengauss, [5.263, 2.582, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]['paly_video'], compare_index_opengauss, [5.263, 2.975, 1, 1], 10, 'red')
            _addText(prs, '120345', compare_index_opengauss, [8.356, 1.968, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["video"], compare_index_opengauss, [8.356, 1.602, 1, 1], 10, 'green')
            _addText(prs, '0', compare_index_opengauss, [8.356, 1.251, 1, 1], 10, 'green')
            _addText(prs, '0', compare_index_opengauss, [10.338, 1.968, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]["video"], compare_index_opengauss, [10.338, 1.602, 1, 1], 10, 'red')
            _addText(prs, '0', compare_index_opengauss, [10.338, 1.251, 1, 1], 10, 'red')
            _addText(prs, data['opengauss']['D0_user'], compare_index_opengauss, [2.877, 3.626, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["contributor"], compare_index_opengauss, [2.877, 3.992, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["committers"], compare_index_opengauss, [2.877, 4.358, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["gitee_star_user"]+data["opengauss"]['github_star'], compare_index_opengauss, [2.877, 4.732, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["gitee_fork_user"]+data["opengauss"]['github_fork'], compare_index_opengauss, [2.877, 5.106, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["gitee_watch_user"]+data["opengauss"]['github_watch'], compare_index_opengauss, [2.877, 5.472, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["gitee_star_user"], compare_index_opengauss, [2.877, 5.838, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["gitee_pr_user"], compare_index_opengauss, [2.877, 6.173, 1, 1], 10, 'green')
            _addText(prs, data["opengauss"]["gitee_issue"], compare_index_opengauss, [2.877, 6.501, 1, 1], 10, 'green')
            _addText(prs, data2['opengauss']['D0_user'], compare_index_opengauss, [5.251, 3.626, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]["contributor"], compare_index_opengauss, [5.251, 3.992, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]["committers"], compare_index_opengauss, [5.251, 4.358, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]["gitee_star_user"]+data2["opengauss"]['github_star'], compare_index_opengauss, [5.251, 4.732, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]["gitee_fork_user"]+data2["opengauss"]['github_fork'], compare_index_opengauss, [5.251, 5.106, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]["gitee_watch_user"]+data2["opengauss"]['github_watch'], compare_index_opengauss, [5.251, 5.472, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]["gitee_star_user"], compare_index_opengauss, [5.251, 5.838, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]["gitee_pr_user"], compare_index_opengauss, [5.251, 6.173, 1, 1], 10, 'red')
            _addText(prs, data2["opengauss"]["gitee_issue"], compare_index_opengauss, [5.251, 6.501, 1, 1], 10, 'red')
        # openlookeng
        home_index_openlookeng = 19
        compare_index_openlookeng = 20
        if not org or org == 'openlookeng':
            _addText(prs, itime, home_index_openlookeng, [4.133, 0.476, 1, 1], 14)  # 触达用户
            _addText(prs, data['openlookeng']['website'], home_index_openlookeng, [5.582, 2.35, 1, 1], 16)  # 触达用户
            _addText(prs, data['openlookeng']['website'], home_index_openlookeng, [7.114, 2.232, 1, 1], 14)  # 网站访问
            _addText(prs, data["openlookeng"]["maillist"], home_index_openlookeng, [8.421, 2.338, 1, 1], 14)  # 邮件列表
            _addText(prs, data["openlookeng"]["repo"], home_index_openlookeng, [9.37, 2.338, 1, 1], 14)  # 代码仓
            _addText(prs, data["openlookeng"]["video"], home_index_openlookeng, [10.472, 2.338, 1, 1], 14)  # 视频
            _addText(prs, data["openlookeng"]["users"], home_index_openlookeng, [4.968, 3.437, 1, 1], 16)  # user
            _addText(prs, data["openlookeng"]["users"], home_index_openlookeng, [7.236, 3.496, 1, 1], 14)  # 试用
            _addText(prs, data["openlookeng"]["contributor"], home_index_openlookeng, [4.645, 4.527, 1, 1], 16)  # Contributor
            _addText(prs, data["openlookeng"]["contributorHuawei"], home_index_openlookeng, [7.318, 4.704, 1, 1], 14)  # 华为
            _addText(prs, data["openlookeng"]["contributorNotHuawei"], home_index_openlookeng, [9.641, 4.633, 1, 1], 14)  # 非华为
            _addText(prs, data["openlookeng"]["committers"], home_index_openlookeng, [4.181, 5.732, 1, 1], 16)  # Committer
            _addText(prs, data["openlookeng"]["committersHuawei"], home_index_openlookeng, [7.377, 5.984, 1, 1], 14)  # 华为
            _addText(prs, data["openlookeng"]["committersNotHuawei"], home_index_openlookeng, [9.685, 6.055, 1, 1], 14)  # 非华为
            # compare
            _addText(prs, data["openlookeng"]['webvisitor'], compare_index_openlookeng, [2.795, 1.098, 1, 1], 10, 'green')
            _addText(prs, data["openlookeng"]['website'], compare_index_openlookeng, [2.795, 1.48, 1, 1], 10, 'green')
            _addText(prs, data["openlookeng"]["repo"], compare_index_openlookeng, [2.795, 1.834, 1, 1], 10, 'green')
            _addText(prs, data["openlookeng"]["maillist"], compare_index_openlookeng, [2.795, 2.208, 1, 1], 10, 'green')
            _addText(prs, data["openlookeng"]["maillistnum"], compare_index_openlookeng, [2.795, 2.582, 1, 1], 10, 'green')
            _addText(prs, data2["openlookeng"]['webvisitor'], compare_index_openlookeng, [5.263, 1.098, 1, 1], 10, 'red')
            _addText(prs, data2["openlookeng"]['website'], compare_index_openlookeng, [5.263, 1.48, 1, 1], 10, 'red')
            _addText(prs, data2["openlookeng"]['repo'], compare_index_openlookeng, [5.263, 1.834, 1, 1], 10, 'red')
            _addText(prs, data2["openlookeng"]['maillist'], compare_index_openlookeng, [5.263, 2.208, 1, 1], 10, 'red')
            _addText(prs, data2["openlookeng"]['maillistnum'], compare_index_openlookeng, [5.263, 2.582, 1, 1], 10, 'red')
            # _addText(prs, '12345', 21, [8.356, 1.968, 1, 1], 10, 'green')
            _addText(prs, data["openlookeng"]["video"], compare_index_openlookeng, [8.356, 1.602, 1, 1], 10, 'green')
            # _addText(prs, '12345', 21, [8.356, 1.251, 1, 1], 10, 'green')
            # _addText(prs, '12345', 21, [10.338, 1.968, 1, 1], 10, 'red')
            _addText(prs, data2["openlookeng"]["video"], compare_index_openlookeng, [10.338, 1.602, 1, 1], 10, 'red')
            # _addText(prs, '12345', 21, [10.338, 1.251, 1, 1], 10, 'red')
            _addText(prs, data["openlookeng"]["contributor"], compare_index_openlookeng, [2.877, 3.992, 1, 1], 10, 'green')
            _addText(prs, data["openlookeng"]["committers"], compare_index_openlookeng, [2.877, 4.358, 1, 1], 10, 'green')
            _addText(prs, data["openlookeng"]["gitee_star_user"], compare_index_openlookeng, [2.877, 4.732, 1, 1], 10,
                     'green')
            _addText(prs, data["openlookeng"]["gitee_fork_user"], compare_index_openlookeng, [2.877, 5.106, 1, 1], 10,
                     'green')
            _addText(prs, data["openlookeng"]["gitee_watch_user"], compare_index_openlookeng, [2.877, 5.472, 1, 1], 10,
                     'green')
            _addText(prs, data["openlookeng"]["gitee_star_user"], compare_index_openlookeng, [2.877, 5.838, 1, 1], 10,
                     'green')
            _addText(prs, data["openlookeng"]["gitee_pr_user"], compare_index_openlookeng, [2.877, 6.173, 1, 1], 10,
                     'green')
            _addText(prs, data2["openlookeng"]["contributor"], compare_index_openlookeng, [5.251, 3.992, 1, 1], 10, 'red')
            _addText(prs, data2["openlookeng"]["committers"], compare_index_openlookeng, [5.251, 4.358, 1, 1], 10, 'red')
            _addText(prs, data2["openlookeng"]["gitee_star_user"], compare_index_openlookeng, [5.251, 4.732, 1, 1], 10,
                     'red')
            _addText(prs, data2["openlookeng"]["gitee_fork_user"], compare_index_openlookeng, [5.251, 5.106, 1, 1], 10,
                     'red')
            _addText(prs, data2["openlookeng"]["gitee_watch_user"], compare_index_openlookeng, [5.251, 5.472, 1, 1], 10,
                     'red')
            _addText(prs, data2["openlookeng"]["gitee_star_user"], compare_index_openlookeng, [5.251, 5.838, 1, 1], 10,
                     'red')
            _addText(prs, data2["openlookeng"]["gitee_pr_user"], compare_index_openlookeng, [5.251, 6.173, 1, 1], 10, 'red')

        prs.save(fileppt)

    def do(self, start='2019-06-01 00:00:00', end=None, org=None, check=None):
        if not end:
            endT = time.strftime("%Y/%m/%d", time.localtime(time.time()))
        else:
            endT = end.split()[0].replace('-', '/')
        startT = start.split()[0].replace('-', '/')
        itime = '%s--%s' % (startT, endT)

        LW_endtime = datetime.datetime.strftime(datetime.datetime.strptime(endT, "%Y/%m/%d") - datetime.timedelta(days=7),
                                                "%Y-%m-%dT%H:%M:%S")
        data = self.getData(self.headers, start, end, org, self.headers)
        LW_data = self.getData(self.headers, start, LW_endtime, org, self.headers)
        data2 = {}
        for key1 in data:
            data2[key1] = {}
            for key2 in data[key1]:
                if data[key1][key2] == 'N/A':
                    data2[key1][key2] = 'N/A'
                else:
                    data2[key1][key2] = int(data[key1][key2]) - int(LW_data[key1][key2])

        fileppt = 'D:\\运营看板报告.pptx'
        fileppt = ''
        # writePPT(data, itime, org, fileppt, data2)

        logo = {'openeuler': 'openEuler', 'mindspore':'MindSpore', 'opengauss':'openGauss', 'openlookeng':'openLooKeng'}
        subject = '开源项目运营看板报告--%s' % logo[org]
        # content = '开源项目运营看板报告：%s' % itime
        content = '''开源项目运营看板报告：%s
    
        内容:
        1.官方网站访问量： %s （较上周上升 %s）
        2.官方网站访客人数： %s （较上周上升 %s）
        3.官方网站视频播放量： %s （较上周上升 %s）
        4.官方网站安装包下载量： %s （较上周上升 %s）
        5.官方网站安装包下载次数： %s （较上周上升 %s）
        6.官方博客数量： %s （较上周上升 %s）
        7.官方博客阅读量： %s （较上周上升 %s）
        8.微信公众号关注人数： %s （较上周上升 %s）
        9.微信社群人数： %s （较上周上升 %s）
        10.微信公众号软文数量： %s （较上周上升 %s）
        11.邮件列表注册人数： %s （较上周上升 %s）
        12.邮件列表邮件数量： %s （较上周上升 %s）
    
        活动：
        1.线上直播次数： %s （较上周上升 %s）
        2.Meetup次数： %s （较上周上升 %s）
        3.B站关注人数： %s （较上周上升 %s）
        4.B站视频播放次数： %s （较上周上升 %s）
    
        开发者：
        1.开发者D0(触达)： %s （较上周上升 %s）
        2.开发者D1(参与)： %s （较上周上升 %s）
        3.开发者D2(贡献)： %s （较上周上升 %s）
        4.Star(Gitee+Github)： %s （较上周上升 %s）
        5.Fork(Gitee+Github)： %s （较上周上升 %s）
        6.Watch(Gitee+Github)： %s （较上周上升 %s）
        7.Gitee组织关注数量： %s （较上周上升 %s）
        8.Pull Request数量： %s （较上周上升 %s）
        9.Issue数量： %s （较上周上升 %s）''' % (itime,
                                       data[org]['website'], data2[org]['website'],
                                       data[org]['webvisitor'], data2[org]['webvisitor'],
                                       data[org]['paly_video'], data2[org]['paly_video'],
                                       data[org]['downlaod_gb'], data2[org]['downlaod_gb'],
                                       data[org]['downlaod'], data2[org]['downlaod'],
                                       'N/A', 'N/A',
                                       'N/A', 'N/A',
                                       'N/A', 'N/A',
                                       'N/A', 'N/A',
                                       'N/A', 'N/A',
                                       data[org]['maillist'], data2[org]['maillist'],
                                       data[org]['maillistnum'], data2[org]['maillistnum'],
                                       # v
                                       'N/A', 'N/A',
                                       'N/A', 'N/A',
                                       'N/A', 'N/A',
                                       # data[org]['video'], data2[org]['video'],
                                       'N/A', 'N/A',
                                       # c
                                       data[org]['downlaod_ip'], data2[org]['downlaod_ip'],
                                       data[org]['contributor'], data2[org]['contributor'],
                                       data[org]['deep_contributor'], data2[org]['deep_contributor'],
                                       data[org]['gitee_star_user'] + data[org]['github_star'],
                                       data2[org]['gitee_star_user'] + data2[org]['github_star'],
                                       data[org]['gitee_fork_user'] + data[org]['github_fork'],
                                       data2[org]['gitee_fork_user'] + data2[org]['github_fork'],
                                       data[org]['gitee_watch_user'] + data[org]['github_watch'],
                                       data2[org]['gitee_watch_user'] + data2[org]['github_watch'],
                                       data[org]['gitee_star_user'], data2[org]['gitee_star_user'],
                                       data[org]['gitee_pr_user'], data2[org]['gitee_pr_user'],
                                       data[org]['gitee_issue'], data2[org]['gitee_issue']
                                       )

        for email in self.email_list:
            send_email('smtp.163.com', 465, self.send, self.smtp, email, subject, content, fileppt)
            time.sleep(6)

