# -*- coding: utf-8 -*-

import bs4
import os
import time
import datetime
from dateutil import tz
import json
import requests
import configparser
import glob
import threading
import requests
import pypistats
import re
from pprint import pprint
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from apscheduler.schedulers.blocking import BlockingScheduler
from data.sendEmail import send_email


def getPng(url, filename, api, start='2020-01-01 00:00:00', end=None):
    starttime = int(round(time.mktime(time.strptime(start, '%Y-%m-%d %H:%M:%S')) * 1000))
    if end:
        endtime = int(round(time.mktime(time.strptime(end, '%Y-%m-%d %H:%M:%S')) * 1000))
    else:
        endtime = int(round(time.time() * 1000))
    time.sleep(5)
    url = re.sub(r'from=\d+&to=\d+&panelId', 'from=%s&to=%s&panelId' % (starttime, endtime), url)
    print(url)
    cmd = 'curl -H "Accept: application/json" -H "Authorization: Bearer %s" "%s" > %s' % (api, url, filename)
    os.system(cmd)


def getData(headers, start='2020-01-01 00:00:00', end=None, org=None, header01=''):
    result = {"openeuler":{}, 'mindspore':{}, 'opengauss':{}, 'openlookeng':{}}
    if end:
        endtime = end.replace(' ', 'T')
    else:
        endtime = time.strftime("%Y-%m-%dT%H:%M:%S", time.localtime(time.time()))
    starttime = start.replace(' ', 'T')
    # openeuler
    data = '''{"size":10000,
  "query": {
    "bool": {
      "filter":{
        "range":{
          "created_at":{
            "gte":"%s.000+0800",
            "lt":"%s.000+0800"
          }
        }
      },
        // "must": [{ "match": { "ip":"127.0.0.1" }}],
        "should": [{
				"match_phrase": {
					"source_type_title": "搜索引擎"
				}
			},
      {
				"match_phrase": {
					"source_type_title": "外部链接"
				}
			},
			{
				"match_phrase": {
					"source_type_title": "直接访问"
				}
			}],
			"minimum_should_match": 1
    }
  },"aggs": {
    "ip_count": {
      "sum": {
        "field": "ip_count"}}}
}''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/baidutongji_openeuler_20200702/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["website"] = int(res["aggregations"]["ip_count"]["value"])

    data = '''{"size":10000,
  "query": {
    "bool": {
      "filter":{
        "range":{
          "created_at":{
            "gte":"%s.000+0800",
            "lt":"%s.000+0800"
          }
        }
      }
    }
  },"aggs": {
    "user_id": {
      "cardinality": {
        "field": "user_id.keyword"}}}
}''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/maillist_user_20200519/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["maillist"] = int(res["aggregations"]["user_id"]["value"])

    data = '''{"size":10000,
  "query": {
    "bool": {
      "filter":{
        "range":{
          "created_at":{
            "gte":"%s.000+0800",
            "lt":"%s.000+0800"
          }
        }
      }, "must": [{ "match": { "event":"DOWNLOAD ZIP" }}]
    }
  },
  "aggs": {
    "author_name": {
      "cardinality": {
        "field": "author_name.keyword"}}}
}''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_log_event_openeuler_all_20200701/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["repo"] = int(res["aggregations"]["author_name"]["value"])

    data = '''{"size":10000,
  "query": {
    "bool": {
      "filter":{
        "range":{
          "created_at":{
            "gte":"%s.000+0800",
            "lt":"%s.000+0800"
          }
        }
      }, "must_not": [{ "match": { "is_topic":1 }}]
    }
  },
  "aggs": {
    "new_followers": {
      "sum": {
        "field": "new_followers"}}}
}''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/bilibili_openeuler_20200710/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["video"] = int(res["aggregations"]["new_followers"]["value"])

    data = '''{"size":10000,
  "query": {
    "bool": {
      "filter":{
        "range":{
          "created_at":{
            "gte":"%s.000+0800",
            "lt":"%s.000+0800"
          }
        }
      }
    }
  },
  "aggs": {
    "sum_value_gb": {
      "sum": {
        "field": "sum_value_gb"}}}
}''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/openeuler_download_20200519/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["users"] = int(res["aggregations"]["sum_value_gb"]["value"] / 6)

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }
        }
      },
      "aggs": {
        "sum_value_gb": {
          "sum": {
            "field": "sum_value_gb"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/openeuler_huaweicloud_vpc/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["downloadall"] = int(res["aggregations"]["sum_value_gb"]["value"] / 6)

    data = '''{"size":10000,
  "query": {
    "bool": {
      "filter":{
        "range":{
          "created_at":{
            "gte":"%s.000+0800",
            "lt":"%s.000+0800"
          }
        }
      },
        // "must": [{ "match": { "ip":"127.0.0.1" }}],
        "should": [{
				"match_phrase": {
					"is_gitee_fork": 1
				}
			},
      {
				"match_phrase": {
					"is_gitee_issue": 1
				}
			},
            {
				"match_phrase": {
					"is_gitee_pull_request": 1
				}
			},
            {
				"match_phrase": {
					"is_gitee_review_comment": 1
				}
			},
            {
				"match_phrase": {
					"is_gitee_comment": 1
				}
			},
			{
				"match_phrase": {
					"is_gitee_issue_comment": 1
				}
			}],
			"minimum_should_match": 1
    }
  },"aggs": {
    "user_login": {
      "cardinality": {
        "field": "user_login.keyword"}}}
}''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_openeuler_all_20200519_2/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["contributor"] = int(res["aggregations"]["user_login"]["value"])

    data = '''{"size":10000,
  "query": {
    "bool": {
      "filter":{
        "range":{
          "created_at":{
            "gte":"%s.000+0800",
            "lt":"%s.000+0800"
          }
        }
      },
        "must": [{ "match": { "is_project_internal_user":1 }}],
        "should": [{
				"match_phrase": {
					"is_gitee_fork": 1
				}
			},
      {
				"match_phrase": {
					"is_gitee_issue": 1
				}
			},
            {
				"match_phrase": {
					"is_gitee_pull_request": 1
				}
			},
            {
				"match_phrase": {
					"is_gitee_review_comment": 1
				}
			},
            {
				"match_phrase": {
					"is_gitee_comment": 1
				}
			},
			{
				"match_phrase": {
					"is_gitee_issue_comment": 1
				}
			}],
			"minimum_should_match": 1
    }
  },"aggs": {
    "user_login": {
      "cardinality": {
        "field": "user_login.keyword"}}}
}''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_openeuler_all_20200519_2/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["contributorHuawei"] = int(res["aggregations"]["user_login"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          },
            "must": [{ "match": { "is_project_internal_user":1 }}],
            "should": [{
    				"match_phrase": {
    					"is_gitee_fork": 1
    				}
    			},
          {
    				"match_phrase": {
    					"is_gitee_issue": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_pull_request": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_review_comment": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_comment": 1
    				}
    			},
    			{
    				"match_phrase": {
    					"is_gitee_issue_comment": 1
    				}
    			}],
    			"minimum_should_match": 1
        }
      },"aggs": {
        "user_login": {
          "cardinality": {
            "field": "user_login.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_openeuler_all_20200519_2/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["contributorNotHuawei"] = int(res["aggregations"]["user_login"]["value"])

    data = '''{"size":10000,
  "query": {
    "bool": {
      "filter":{
        "range":{
          "created_at":{
            "gte":"%s.000+0800",
            "lt":"%s.000+0800"
          }
        }
      }
    }
  },"aggs": {
    "user_gitee_name": {
      "cardinality": {
        "field": "user_gitee_name.keyword"}}}
}''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/openeuler_sig_20200629/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["committers"] = int(res["aggregations"]["user_gitee_name"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          },"must": [{ "match": { "is_project_internal_user":1 }}]
        }
      },"aggs": {
        "user_gitee_name": {
          "cardinality": {
            "field": "user_gitee_name.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/openeuler_sig_20200629/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["committersHuawei"] = int(res["aggregations"]["user_gitee_name"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          },"must_not": [{ "match": { "is_project_internal_user":1 }}]
        }
      },"aggs": {
        "user_gitee_name": {
          "cardinality": {
            "field": "user_gitee_name.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/openeuler_sig_20200629/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["committersNotHuawei"] = int(res["aggregations"]["user_gitee_name"]["value"])

    # mindspore
    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          },
            // "must": [{ "match": { "ip":"127.0.0.1" }}],
            "should": [{
    				"match_phrase": {
    					"source_type_title": "搜索引擎"
    				}
    			},
          {
    				"match_phrase": {
    					"source_type_title": "外部链接"
    				}
    			},
    			{
    				"match_phrase": {
    					"source_type_title": "直接访问"
    				}
    			}],
    			"minimum_should_match": 1
        }
      },"aggs": {
        "ip_count": {
          "sum": {
            "field": "ip_count"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/baidutongji_mindspore_20200603/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["mindspore"]["website"] = int(res["aggregations"]["ip_count"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }
        }
      },"aggs": {
        "user_id": {
          "cardinality": {
            "field": "user_id.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://159.138.2.100:9200/mindspore_hyperkitty_enriched/_search'
    res = json.loads(requests.get(url=url, headers=header01, verify=False, data=data.encode('utf-8')).content)
    result["mindspore"]["maillist"] = int(res["aggregations"]["user_id"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }, "must": [{ "match": { "event":"DOWNLOAD ZIP" }},{ "match": { "is_forked_repo":0 }}]
        }
      },
      "aggs": {
        "author_name": {
          "cardinality": {
            "field": "author_name.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_log_event_mindspore_all_20200604/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["mindspore"]["repo"] = int(res["aggregations"]["author_name"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }, "must_not": [{ "match": { "is_topic":1 }}]
        }
      },
      "aggs": {
        "new_followers": {
          "sum": {
            "field": "new_followers"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/bilibili_mindspore_20200710/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["mindspore"]["video"] = int(res["aggregations"]["new_followers"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }
        }
      },
      "aggs": {
        "ip": {
          "cardinality": {
            "field": "ip.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/mindspore_obs_20200515/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["mindspore"]["users"] = int(res["aggregations"]["ip"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          },
            // "must": [{ "match": { "ip":"127.0.0.1" }}],
            "should": [{
    				"match_phrase": {
    					"is_gitee_fork": 1
    				}
    			},
          {
    				"match_phrase": {
    					"is_gitee_issue": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_pull_request": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_review_comment": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_comment": 1
    				}
    			},
    			{
    				"match_phrase": {
    					"is_gitee_issue_comment": 1
    				}
    			}],
    			"minimum_should_match": 1
        }
      },"aggs": {
        "user_login": {
          "cardinality": {
            "field": "user_login.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_test_mindspore_all_20200511/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["mindspore"]["contributor"] = int(res["aggregations"]["user_login"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          },
            "must": [{ "match": { "is_project_internal_user":1 }}],
            "should": [{
    				"match_phrase": {
    					"is_gitee_fork": 1
    				}
    			},
          {
    				"match_phrase": {
    					"is_gitee_issue": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_pull_request": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_review_comment": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_comment": 1
    				}
    			},
    			{
    				"match_phrase": {
    					"is_gitee_issue_comment": 1
    				}
    			}],
    			"minimum_should_match": 1
        }
      },"aggs": {
        "user_login": {
          "cardinality": {
            "field": "user_login.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_test_mindspore_all_20200511/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["mindspore"]["contributorHuawei"] = int(res["aggregations"]["user_login"]["value"])

    data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              },
                "must": [{ "match": { "is_project_internal_user":1 }}],
                "should": [{
        				"match_phrase": {
        					"is_gitee_fork": 1
        				}
        			},
              {
        				"match_phrase": {
        					"is_gitee_issue": 1
        				}
        			},
                    {
        				"match_phrase": {
        					"is_gitee_pull_request": 1
        				}
        			},
                    {
        				"match_phrase": {
        					"is_gitee_review_comment": 1
        				}
        			},
                    {
        				"match_phrase": {
        					"is_gitee_comment": 1
        				}
        			},
        			{
        				"match_phrase": {
        					"is_gitee_issue_comment": 1
        				}
        			}],
        			"minimum_should_match": 1
            }
          },"aggs": {
            "user_login": {
              "cardinality": {
                "field": "user_login.keyword"}}}
        }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_test_mindspore_all_20200511/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["contributorNotHuawei"] = int(res["aggregations"]["user_login"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          },"must": [{ "match": { "is_committer":1 }}]
        }
      },"aggs": {
        "user_login": {
          "cardinality": {
            "field": "user_login.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_test_mindspore_all_20200511/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["mindspore"]["committers"] = int(res["aggregations"]["user_login"]["value"])

    # opengauss
    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          },
            // "must": [{ "match": { "ip":"127.0.0.1" }}],
            "should": [{
    				"match_phrase": {
    					"source_type_title": "搜索引擎"
    				}
    			},
          {
    				"match_phrase": {
    					"source_type_title": "外部链接"
    				}
    			},
    			{
    				"match_phrase": {
    					"source_type_title": "直接访问"
    				}
    			}],
    			"minimum_should_match": 1
        }
      },"aggs": {
        "ip_count": {
          "sum": {
            "field": "ip_count"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/baidutongji_opengauss_20200702/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["opengauss"]["website"] = int(res["aggregations"]["ip_count"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }
        }
      },"aggs": {
        "user_id": {
          "cardinality": {
            "field": "user_id.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/opengauss_maillist_user_20200609/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["opengauss"]["maillist"] = int(res["aggregations"]["user_id"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }, "must": [{ "match": { "event":"DOWNLOAD ZIP" }},{ "match": { "is_forked_repo":1 }}]
        }
      },
      "aggs": {
        "author_name": {
          "cardinality": {
            "field": "author_name.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_log_event_opengauss_all_20200703/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["opengauss"]["repo"] = int(res["aggregations"]["author_name"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }, "must_not": [{ "match": { "is_topic":1 }}]
        }
      },
      "aggs": {
        "new_followers": {
          "sum": {
            "field": "new_followers"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/bilibili_opengauss_20200710/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["opengauss"]["video"] = int(res["aggregations"]["new_followers"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }
        }
      },
      "aggs": {
        "ip": {
          "cardinality": {
            "field": "ip.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/opengauss_obs_20200708/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["opengauss"]["users"] = int(res["aggregations"]["ip"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          },
            // "must": [{ "match": { "ip":"127.0.0.1" }}],
            "should": [{
    				"match_phrase": {
    					"is_gitee_fork": 1
    				}
    			},
          {
    				"match_phrase": {
    					"is_gitee_issue": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_pull_request": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_review_comment": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_comment": 1
    				}
    			},
    			{
    				"match_phrase": {
    					"is_gitee_issue_comment": 1
    				}
    			}],
    			"minimum_should_match": 1
        }
      },"aggs": {
        "user_login": {
          "cardinality": {
            "field": "user_login.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_opengauss_all_20200513/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["opengauss"]["contributor"] = int(res["aggregations"]["user_login"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          },
            "must": [{ "match": { "is_project_internal_user":1 }}],
            "should": [{
    				"match_phrase": {
    					"is_gitee_fork": 1
    				}
    			},
          {
    				"match_phrase": {
    					"is_gitee_issue": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_pull_request": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_review_comment": 1
    				}
    			},
                {
    				"match_phrase": {
    					"is_gitee_comment": 1
    				}
    			},
    			{
    				"match_phrase": {
    					"is_gitee_issue_comment": 1
    				}
    			}],
    			"minimum_should_match": 1
        }
      },"aggs": {
        "user_login": {
          "cardinality": {
            "field": "user_login.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_opengauss_all_20200513/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["opengauss"]["contributorHuawei"] = int(res["aggregations"]["user_login"]["value"])

    data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              },
                "must": [{ "match": { "is_project_internal_user":1 }}],
                "should": [{
        				"match_phrase": {
        					"is_gitee_fork": 1
        				}
        			},
              {
        				"match_phrase": {
        					"is_gitee_issue": 1
        				}
        			},
                    {
        				"match_phrase": {
        					"is_gitee_pull_request": 1
        				}
        			},
                    {
        				"match_phrase": {
        					"is_gitee_review_comment": 1
        				}
        			},
                    {
        				"match_phrase": {
        					"is_gitee_comment": 1
        				}
        			},
        			{
        				"match_phrase": {
        					"is_gitee_issue_comment": 1
        				}
        			}],
        			"minimum_should_match": 1
            }
          },"aggs": {
            "user_login": {
              "cardinality": {
                "field": "user_login.keyword"}}}
        }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_opengauss_all_20200513/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["opengauss"]["contributorNotHuawei"] = int(res["aggregations"]["user_login"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }
        }
      },"aggs": {
        "user_login": {
          "cardinality": {
            "field": "user_login.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/opengauss_sig_20200716/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["opengauss"]["committers"] = int(res["aggregations"]["user_login"]["value"])

    data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              },"must": [{ "match": { "is_project_internal_user":1 }}]
            }
          },"aggs": {
            "user_login": {
              "cardinality": {
                "field": "user_login.keyword"}}}
        }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/opengauss_sig_20200716/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openeuler"]["committersHuawei"] = int(res["aggregations"]["user_login"]["value"])

    data = '''{"size":10000,
          "query": {
            "bool": {
              "filter":{
                "range":{
                  "created_at":{
                    "gte":"%s.000+0800",
                    "lt":"%s.000+0800"
                  }
                }
              },"must_not": [{ "match": { "is_project_internal_user":1 }}]
            }
          },"aggs": {
            "user_gitee_name": {
              "cardinality": {
                "field": "user_gitee_name.keyword"}}}
        }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/opengauss_sig_20200716/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["opengauss"]["committersNotHuawei"] = int(res["aggregations"]["user_gitee_name"]["value"])

    # openlookeng
    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          },
            // "must": [{ "match": { "ip":"127.0.0.1" }}],
            "should": [{
    				"match_phrase": {
    					"source_type_title": "搜索引擎"
    				}
    			},
          {
    				"match_phrase": {
    					"source_type_title": "外部链接"
    				}
    			},
    			{
    				"match_phrase": {
    					"source_type_title": "直接访问"
    				}
    			}],
    			"minimum_should_match": 1
        }
      },"aggs": {
        "ip_count": {
          "sum": {
            "field": "ip_count"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/baidutongji_openlookeng_20200629/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openlookeng"]["website"] = int(res["aggregations"]["ip_count"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }
        }
      },"aggs": {
        "user_id": {
          "cardinality": {
            "field": "user_id.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/maillist_user_openlookeng_20200702/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openlookeng"]["maillist"] = int(res["aggregations"]["user_id"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }, "must": [{ "match": { "event":"DOWNLOAD ZIP" }},{ "match": { "is_forked_repo":1 }}]
        }
      },
      "aggs": {
        "author_name": {
          "cardinality": {
            "field": "author_name.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_log_event_openlookeng_all_20200703/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openlookeng"]["repo"] = int(res["aggregations"]["author_name"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }, "must_not": [{ "match": { "is_topic":1 }}]
        }
      },
      "aggs": {
        "new_followers": {
          "sum": {
            "field": "new_followers"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/bilibili_openlookeng_20200710/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openlookeng"]["video"] = int(res["aggregations"]["new_followers"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }, "must_not": [{ "match": { "path.keyword":"-" }}]
        }
      },
      "aggs": {
        "ip": {
          "cardinality": {
            "field": "ip.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/openlookeng_obs_20200706/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["openlookeng"]["users"] = int(res["aggregations"]["ip"]["value"])

    data = '''{"size":10000,
      "query": {
        "bool": {
          "filter":{
            "range":{
              "created_at":{
                "gte":"%s.000+0800",
                "lt":"%s.000+0800"
              }
            }
          }, "must": [{ "match": { "is_first_contribute":1 }}]
        }
      },"aggs": {
        "user_login": {
          "cardinality": {
            "field": "user_login.keyword"}}}
    }''' % (starttime, endtime)
    url = 'https://119.8.111.61:9200/gitee_openlookeng_all_20200629/_search'
    res = json.loads(requests.get(url=url, headers=headers, verify=False, data=data.encode('utf-8')).content)
    result["opengauss"]["contributor"] = int(res["aggregations"]["user_login"]["value"])

    return result


def writePPT(data, itime, org=None, fileppt=''):
    def _addText(prs, content, slides, inches, size):
        textbox = prs.slides[slides].shapes.add_textbox(Inches(inches[0]), Inches(inches[1]), Inches(inches[2]),
                                                        Inches(inches[3]))
        tf = textbox.text_frame
        para = tf.add_paragraph()
        para.text = str(content)
        font = para.font
        font.size = Pt(size)

    prs = Presentation('模板.pptx')

    # add picture
    # openeuler
    if not org or org == 'openeuler':
        prs.slides[2].shapes.add_picture('D:\\GrafanaPng\\openeuler_webviews.png', Inches(1), Inches(1.5), Inches(11.3),
                                         Inches(3))
        prs.slides[2].shapes.add_picture('D:\\GrafanaPng\\openeuler_maillist.png', Inches(1), Inches(4.72), Inches(6),
                                         Inches(2))
        prs.slides[2].shapes.add_picture('D:\\GrafanaPng\\openeuler_repo.png', Inches(7.1), Inches(4.72), Inches(6),
                                         Inches(2))
        prs.slides[3].shapes.add_picture('D:\\GrafanaPng\\openeuler_users.png', Inches(1), Inches(1.5), Inches(11.3),
                                         Inches(5))
        prs.slides[4].shapes.add_picture('D:\\GrafanaPng\\openeuler_contributorlines.png', Inches(1), Inches(1.5),
                                         Inches(11.3), Inches(5))
        prs.slides[5].shapes.add_picture('D:\\GrafanaPng\\openeuler_committerlines.png', Inches(1), Inches(1.5),
                                         Inches(11.3), Inches(5))

    # mindspore
    if not org or org == 'mindspore':
        prs.slides[7].shapes.add_picture('D:\\GrafanaPng\\mindspore_webviews.png', Inches(1), Inches(1.5), Inches(11.3),
                                         Inches(3))
        prs.slides[7].shapes.add_picture('D:\\GrafanaPng\\mindspore_maillist.png', Inches(1), Inches(4.72), Inches(6),
                                         Inches(2))
        prs.slides[7].shapes.add_picture('D:\\GrafanaPng\\mindspore_repo.png', Inches(7.1), Inches(4.72), Inches(6),
                                         Inches(2))
        prs.slides[8].shapes.add_picture('D:\\GrafanaPng\\mindspore_users.png', Inches(1), Inches(1.5), Inches(11.3),
                                         Inches(5))
        prs.slides[9].shapes.add_picture('D:\\GrafanaPng\\mindspore_contributorlines.png', Inches(1), Inches(1.5),
                                         Inches(11.3), Inches(5))
        prs.slides[10].shapes.add_picture('D:\\GrafanaPng\\mindspore_committerlines.png', Inches(1), Inches(1.5),
                                          Inches(11.3), Inches(5))

    # opengauss
    if not org or org == 'opengauss':
        prs.slides[12].shapes.add_picture('D:\\GrafanaPng\\opengauss_webviews.png', Inches(1), Inches(1.5), Inches(11.3),
                                          Inches(3))
        prs.slides[12].shapes.add_picture('D:\\GrafanaPng\\opengauss_maillist.png', Inches(1), Inches(4.72), Inches(6),
                                          Inches(2))
        prs.slides[12].shapes.add_picture('D:\\GrafanaPng\\opengauss_repo.png', Inches(7.1), Inches(4.72), Inches(6),
                                          Inches(2))
        prs.slides[13].shapes.add_picture('D:\\GrafanaPng\\opengauss_users.png', Inches(1), Inches(1.5), Inches(11.3),
                                          Inches(5))
        prs.slides[14].shapes.add_picture('D:\\GrafanaPng\\opengauss_contributorlines.png', Inches(1), Inches(1.5),
                                          Inches(11.3), Inches(5))
        prs.slides[15].shapes.add_picture('D:\\GrafanaPng\\opengauss_committerlines.png', Inches(1), Inches(1.5),
                                          Inches(11.3), Inches(5))

    # openlookeng
    if not org or org == 'openlookeng':
        prs.slides[17].shapes.add_picture('D:\\GrafanaPng\\openlookeng_webviews.png', Inches(1), Inches(1.5), Inches(11.3),
                                          Inches(3))
        prs.slides[17].shapes.add_picture('D:\\GrafanaPng\\openlookeng_maillist.png', Inches(1), Inches(4.72), Inches(6),
                                          Inches(2))
        prs.slides[17].shapes.add_picture('D:\\GrafanaPng\\openlookeng_repo.png', Inches(7.1), Inches(4.72), Inches(6),
                                          Inches(2))
        prs.slides[18].shapes.add_picture('D:\\GrafanaPng\\openlookeng_users.png', Inches(1), Inches(1.5), Inches(11.3),
                                          Inches(5))
        prs.slides[19].shapes.add_picture('D:\\GrafanaPng\\openlookeng_contributorlines.png', Inches(1), Inches(1.5),
                                          Inches(11.3), Inches(5))
        # prs.slides[20].shapes.add_picture('D:\\GrafanaPng\\openlookeng_committerlines.png', Inches(1), Inches(1.5),
        #                                   Inches(11.3), Inches(5))

    # add text
    # openeuler
    if not org or org == 'openeuler':
        _addText(prs, itime, 1, [3.397, 0.476, 1, 1], 14)
        _addText(prs, data['openeuler']['website'], 1, [5.582, 2.35, 1, 1], 16)  # 触达用户
        _addText(prs, data['openeuler']['website'], 1, [7.114, 2.232, 1, 1], 14)  # 网站访问
        _addText(prs, data["openeuler"]["maillist"], 1, [8.421, 2.338, 1, 1], 14)  # 邮件列表
        _addText(prs, data["openeuler"]["repo"], 1, [9.37, 2.338, 1, 1], 14)  # 代码仓
        _addText(prs, data["openeuler"]["video"], 1, [10.472, 2.338, 1, 1], 14)  # 视频
        _addText(prs, data["openeuler"]["users"], 1, [4.968, 3.437, 1, 1], 16)  # user
        _addText(prs, data["openeuler"]["downloadall"], 1, [7.236, 3.496, 1, 1], 14)  # 试用
        _addText(prs, data["openeuler"]["contributor"], 1, [4.645, 4.527, 1, 1], 16)  # Contributor
        _addText(prs, data["openeuler"]["contributorHuawei"], 1, [7.318, 4.704, 1, 1], 14)  # 华为
        _addText(prs, data["openeuler"]["contributorNotHuawei"], 1, [9.641, 4.633, 1, 1], 14)  # 非华为
        _addText(prs, data["openeuler"]["committers"], 1, [4.181, 5.732, 1, 1], 16)  # Committer
        _addText(prs, data["openeuler"]["committersHuawei"], 1, [7.377, 5.984, 1, 1], 14)  # 华为
        _addText(prs, data["openeuler"]["committersNotHuawei"], 1, [9.685, 6.055, 1, 1], 14)  # 非华为
    # mindspore
    if not org or org == 'mindspore':
        _addText(prs, itime, 6, [3.397, 0.476, 1, 1], 14)
        _addText(prs, data['mindspore']['website'], 6, [5.582, 2.35, 1, 1], 16)  # 触达用户
        _addText(prs, data['mindspore']['website'], 6, [7.114, 2.232, 1, 1], 14)  # 网站访问
        _addText(prs, data["mindspore"]["maillist"], 6, [8.421, 2.338, 1, 1], 14)  # 邮件列表
        _addText(prs, data["mindspore"]["repo"], 6, [9.37, 2.338, 1, 1], 14)  # 代码仓
        _addText(prs, data["mindspore"]["video"], 6, [10.472, 2.338, 1, 1], 14)  # 视频
        _addText(prs, data["mindspore"]["users"], 6, [4.968, 3.437, 1, 1], 16)  # user
        _addText(prs, data["mindspore"]["users"], 6, [7.236, 3.496, 1, 1], 14)  # 试用
        _addText(prs, data["mindspore"]["contributor"], 6, [4.645, 4.527, 1, 1], 16)  # Contributor
        _addText(prs, data["mindspore"]["contributorHuawei"], 6, [7.318, 4.704, 1, 1], 14)  # 华为
        _addText(prs, data["mindspore"]["contributorNotHuawei"], 6, [9.641, 4.633, 1, 1], 14)  # 非华为
        _addText(prs, data["mindspore"]["committers"], 6, [4.181, 5.732, 1, 1], 16)  # Committer
        _addText(prs, data["mindspore"]["committersHuawei"], 6, [7.377, 5.984, 1, 1], 14)  # 华为
        _addText(prs, data["mindspore"]["committersNotHuawei"], 6, [9.685, 6.055, 1, 1], 14)  # 非华为
    # opengauss
    if not org or org == 'opengauss':
        _addText(prs, itime, 11, [3.397, 0.476, 1, 1], 14)  # 触达用户
        _addText(prs, data['opengauss']['website'], 11, [5.582, 2.35, 1, 1], 16)  # 触达用户
        _addText(prs, data['opengauss']['website'], 11, [7.114, 2.232, 1, 1], 14)  # 网站访问
        _addText(prs, data["opengauss"]["maillist"], 11, [8.421, 2.338, 1, 1], 14)  # 邮件列表
        _addText(prs, data["opengauss"]["repo"], 11, [9.37, 2.338, 1, 1], 14)  # 代码仓
        _addText(prs, data["opengauss"]["video"], 11, [10.472, 2.338, 1, 1], 14)  # 视频
        _addText(prs, data["opengauss"]["users"], 11, [4.968, 3.437, 1, 1], 16)  # user
        _addText(prs, data["opengauss"]["users"], 11, [7.236, 3.496, 1, 1], 14)  # 试用
        _addText(prs, data["opengauss"]["contributor"], 11, [4.645, 4.527, 1, 1], 16)  # Contributor
        _addText(prs, data["opengauss"]["contributorHuawei"], 11, [7.318, 4.704, 1, 1], 14)  # 华为
        _addText(prs, data["opengauss"]["contributorNotHuawei"], 11, [9.641, 4.633, 1, 1], 14)  # 非华为
        _addText(prs, data["opengauss"]["committers"], 11, [4.181, 5.732, 1, 1], 16)  # Committer
        _addText(prs, data["opengauss"]["committersHuawei"], 11, [7.377, 5.984, 1, 1], 14)  # 华为
        _addText(prs, data["opengauss"]["committersNotHuawei"], 11, [9.685, 6.055, 1, 1], 14)  # 非华为
    # openlookeng
    if not org or org == 'openlookeng':
        _addText(prs, itime, 16, [4.133, 0.476, 1, 1], 14)  # 触达用户
        _addText(prs, data['openlookeng']['website'], 16, [5.582, 2.35, 1, 1], 16)  # 触达用户
        _addText(prs, data['openlookeng']['website'], 16, [7.114, 2.232, 1, 1], 14)  # 网站访问
        _addText(prs, data["openlookeng"]["maillist"], 16, [8.421, 2.338, 1, 1], 14)  # 邮件列表
        _addText(prs, data["openlookeng"]["repo"], 16, [9.37, 2.338, 1, 1], 14)  # 代码仓
        _addText(prs, data["openlookeng"]["video"], 16, [10.472, 2.338, 1, 1], 14)  # 视频
        _addText(prs, data["openlookeng"]["users"], 16, [4.968, 3.437, 1, 1], 16)  # user
        _addText(prs, data["openlookeng"]["users"], 16, [7.236, 3.496, 1, 1], 14)  # 试用
        _addText(prs, data["openlookeng"]["contributor"], 16, [4.645, 4.527, 1, 1], 16)  # Contributor
        _addText(prs, data["openlookeng"]["contributorHuawei"], 16, [7.318, 4.704, 1, 1], 14)  # 华为
        _addText(prs, data["openlookeng"]["contributorNotHuawei"], 16, [9.641, 4.633, 1, 1], 14)  # 非华为
        _addText(prs, data["openlookeng"]["committers"], 16, [4.181, 5.732, 1, 1], 16)  # Committer
        _addText(prs, data["openlookeng"]["committersHuawei"], 16, [7.377, 5.984, 1, 1], 14)  # 华为
        _addText(prs, data["openlookeng"]["committersNotHuawei"], 16, [9.685, 6.055, 1, 1], 14)  # 非华为

    prs.save(fileppt)


def run(start='2020-01-01 00:00:00', end=None, org=None):
    if not os.path.exists('D:\\GrafanaPng'):
        os.makedirs('D:\\GrafanaPng')

    with open("pngurl.json", 'r') as fj:
        f = json.load(fj)
    if org:
        for c in f[org]['pngurl']:
            if f[org]['pngurl'][c]:
                filename = "D:\\GrafanaPng\\%s_%s.png" % (org, c)
                # getPng(f[org]['pngurl'][c], filename, f[org]['APIkey'], start, end)
    else:
        for z in f:
            if 'pngurl' in f[z]:
                for c in f[z]['pngurl']:
                    if f[z]['pngurl'][c]:
                        filename = "D:\\GrafanaPng\\%s_%s.png" % (z, c)
                        # getPng(f[z]['pngurl'][c], filename, f[z]['APIkey'], start, end)

    if not end:
        endT = time.strftime("%Y/%m/%d", time.localtime(time.time()))
    else:
        endT = end.split()[0].replace('-', '/')
    startT = start.split()[0].replace('-', '/')
    itime = '%s--%s' % (startT, endT)

    header = f['openeuler']['header']
    header01 = f['openeuler']['header01']
    data = getData(header, start, end, org, header01)
    fileppt = 'D:\\运营看板报告.pptx'
    writePPT(data, itime, org, fileppt)

    subject = '开源项目运营看板报告'
    content = '开源项目运营看板报告：%s' % itime
    for email in f['emaillist']:
        send_email('smtp.163.com', 465, 'grafanareport@163.com', 'XTMTJVUUDQZRWPKG', email, subject, content, fileppt)
        time.sleep(6)


if __name__ == "__main__":


    # scheduler = BlockingScheduler()
    # # 没5小时run一次
    # scheduler.add_job(run, trigger="interval", hour=5, args=())
    #
    # scheduler.start()


    start = '2020-07-01 00:00:00'
    end = '2020-07-31 00:00:00'
    run(start=start, end=end, org='openeuler')


